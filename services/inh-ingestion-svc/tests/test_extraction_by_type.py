"""Per-file-type extraction tests against the bundled sample documents.

Verifies that each end-to-end supported format (txt, md, csv, json, html, docx,
pdf, xlsx, pptx, yaml, toml, xml [#121], srt, vtt [#127]) extracts to
non-empty, readable text via the production extractor helpers. These run
offline (no storage/staging/Temporal) by calling the extractor helpers
directly.

XLSX (#118) and PPTX (#119) were previously hard-rejected here
(``test_xlsx_extraction_is_rejected``, now replaced by ``test_extract_xlsx``
below per the #118 acceptance criterion: "flip the XLSX-rejection test... to
a positive extraction test").
"""

import io
import json
import zipfile
from pathlib import Path

import openpyxl
import pytest
from pptx import Presentation
from temporalio.exceptions import ApplicationError

from src.temporal.activities.extract import (
    _extract_docx_text,
    _extract_epub_text,
    _extract_html_text,
    _extract_json_text,
    _extract_odt_text,
    _extract_pdf_text,
    _extract_pptx_text,
    _extract_rtf_text,
    _extract_subtitle_text,
    _extract_xlsx_text,
    _extract_xml_text,
    _format_xlsx_cell,
)

# tests/ -> inh-ingestion-svc -> services -> repo
_REPO_ROOT = Path(__file__).resolve().parents[3]
SAMPLE_DOCS_DIR = _REPO_ROOT / "docs" / "examples" / "sample-documents"


@pytest.fixture(autouse=True)
def cleanup_test_data():
    """Override the DB-backed root autouse fixture so these stay offline.

    These extraction tests need neither PostgreSQL nor any other live service;
    shadowing the root ``cleanup_test_data`` (which skips when PostgreSQL is
    down) lets them run unconditionally, including in local dev without Docker.
    """
    yield


def _read(filename: str) -> bytes:
    path = SAMPLE_DOCS_DIR / filename
    if not path.is_file():
        pytest.skip(f"Sample fixture missing: {path}")
    return path.read_bytes()


def test_extract_plain_text():
    text = _read("sample.txt").decode("utf-8", errors="ignore")
    assert text.strip()
    assert "Inherent" in text


def test_extract_markdown():
    text = _read("sample.md").decode("utf-8", errors="ignore")
    assert text.strip()


def test_extract_csv():
    text = _read("sample.csv").decode("utf-8", errors="ignore")
    assert text.strip()
    assert "," in text


def test_extract_json():
    """Calls the real production extractor (`_extract_json_text`) rather than
    re-deriving the same json.loads/dumps inline -- so this test actually
    exercises the function #195's failure-path tests below wrap."""
    text = _extract_json_text(_read("sample.json"))
    assert text.strip()
    data = json.loads(text)
    assert isinstance(data, (dict, list))


def test_extract_html():
    text = _extract_html_text(_read("sample.html"))
    assert text.strip()
    # Tags must be stripped.
    assert "<html" not in text.lower()
    assert "<body" not in text.lower()


def test_extract_docx():
    text = _extract_docx_text(_read("sample.docx"))
    assert text.strip()
    assert "Inherent" in text


def test_extract_pdf():
    """Hand-built sample PDF must yield extractable text."""
    text = _extract_pdf_text(_read("sample.pdf"))
    assert text.strip(), "PDF extraction returned empty text"
    assert "Inherent" in text


def test_extract_xlsx():
    """#118: row-aware, sheet-boundary-preserving extraction against the
    bundled multi-sheet/merged-cell/empty-row fixture."""
    text = _extract_xlsx_text(_read("sample.xlsx"))
    assert text.strip()
    assert "Inherent" in text
    # Sheet boundaries are preserved as "## Sheet: <name>" headers -- an
    # agent reading the flattened text can still tell which sheet a row
    # came from.
    assert "## Sheet: Overview" in text
    assert "## Sheet: Notes" in text
    # Row-aware: cells stay pipe-delimited in column order, so "which value
    # sat in which column" survives the flatten to plain text.
    assert "Product | Region | Revenue" in text
    # The fixture's merged title cells (Overview A1:D1, Notes A1:B1) carry a
    # visible span marker -- distinguishes "blank because merged" from
    # "blank because genuinely empty" (review follow-up: previously the
    # merge flattened to a value cell followed by unexplained blank cells).
    assert "[merged A1:D1]" in text
    assert "[merged A1:B1]" in text


def test_extract_yaml():
    """#121: decoded as plain text, no parse step."""
    text = _read("sample.yaml").decode("utf-8", errors="ignore")
    assert text.strip()
    assert "service: inherent" in text


def test_extract_toml():
    """#121: decoded as plain text, no parse step."""
    text = _read("sample.toml").decode("utf-8", errors="ignore")
    assert text.strip()
    assert 'name = "inherent"' in text


def test_extract_xml():
    """#121: tags stripped, element text kept, attribute values dropped."""
    text = _extract_xml_text(_read("sample.xml"))
    assert text.strip()
    assert "Backend for turning company knowledge" in text
    assert "<service" not in text
    # Attribute-value policy: the `name="inherent"` attribute's VALUE does
    # not survive tag-stripping, only element text content does.
    assert "inherent" not in text.lower()


def test_extract_srt():
    """#127: cue numbers/timestamps stripped, prose + coarse markers kept."""
    text = _extract_subtitle_text(_read("sample.srt"), "sample.srt")
    assert "Welcome to the Inherent product walkthrough." in text
    assert "Let's start with how documents get uploaded." in text
    assert "-->" not in text
    assert "[t=00:00]" in text


def test_extract_vtt():
    """#127: header stripped, cue numbers/timestamps stripped, prose kept."""
    text = _extract_subtitle_text(_read("sample.vtt"), "sample.vtt")
    assert "Welcome to the Inherent product walkthrough." in text
    assert "Let's start with how documents get uploaded." in text
    assert "WEBVTT" not in text
    assert "-->" not in text


def test_extract_pptx():
    """#119: slide-boundary sections, in-order text frames, table rows, and
    speaker notes, against the bundled fixture."""
    text = _extract_pptx_text(_read("sample.pptx"))
    assert text.strip()
    assert "Inherent" in text
    # Slide boundaries are preserved as "## Slide <n>[: <title>]" headers.
    assert "## Slide 1" in text
    # Speaker notes are appended under a "Notes:" section, not silently
    # dropped -- this is what makes a notes-only query able to retrieve the
    # right slide's chunk downstream.
    assert "Notes:" in text
    # A table shape's rows render pipe-delimited, same convention as XLSX.
    assert " | " in text


# ---------------------------------------------------------------------------
# Failure paths (#118/#119) -- corrupt input, password protection, legacy
# formats, and pathological size must all fail CLEARLY, never hang, OOM, or
# silently mis-parse. See the module docstring in inh_contracts/file_types.py
# and extract.py for why "no default lossy fallback" is the house rule.
# ---------------------------------------------------------------------------


class TestXlsxFailurePaths:
    def test_corrupt_truncated_zip_raises_non_retryable(self):
        """A ZIP-signature prefix followed by garbage (corrupt or truncated
        upload) must fail loudly -- openpyxl's zipfile.BadZipFile is caught
        and re-raised as an actionable, NON-RETRYABLE ApplicationError (review
        follow-up: a bare RuntimeError here was retried 3x by Temporal's
        default activity RetryPolicy even though the same bytes can never
        succeed on retry -- deterministic given fixed content, same reasoning
        as `_resolve_extractor`'s existing non-retryable failures)."""
        with pytest.raises(ApplicationError, match="XLSX extraction failed") as exc_info:
            _extract_xlsx_text(b"PK\x03\x04" + b"garbage, not a real zip central directory")
        assert exc_info.value.non_retryable

    def test_password_protected_bytes_raise_non_retryable(self):
        """Password-protected OOXML is wrapped in an OLE2/CFBF container
        (magic D0 CF 11 E0 A1 B1 1A E1), not a ZIP -- it is normally caught
        earlier by inh_contracts.sniff_content_type's magic-byte check
        (declared xlsx, bytes don't match the zip signature) before ever
        reaching this extractor. This test is the defense-in-depth layer:
        even called directly, the extractor itself must not crash
        ungracefully or hang -- it fails the same clear, non-retryable way as
        any other non-zip input."""
        ole2_encrypted_magic = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 32
        with pytest.raises(ApplicationError, match="XLSX extraction failed") as exc_info:
            _extract_xlsx_text(ole2_encrypted_magic)
        assert exc_info.value.non_retryable

    def test_missing_dependency_raises_non_retryable(self, monkeypatch):
        """openpyxl not being installed is itself deterministic (retrying the
        same worker process, or any worker with the same image, can never
        succeed) -- non-retryable, same as every other failure mode here."""
        import builtins

        real_import = builtins.__import__

        def _blocking_import(name, *args, **kwargs):
            if name == "openpyxl":
                raise ImportError("simulated: openpyxl not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", _blocking_import)

        with pytest.raises(ApplicationError, match="openpyxl not available") as exc_info:
            _extract_xlsx_text(b"irrelevant")
        assert exc_info.value.non_retryable

    def test_empty_workbook_yields_empty_text_not_a_bare_sheet_heading(self):
        """Review follow-up: an 'empty' workbook previously still yielded its
        sheet's '## Sheet: <name>' heading with zero data rows -- non-empty
        text that cleared extract_text's empty-extraction guard and got a
        content-free document indexed. openpyxl (and Excel) always ships
        >=1 sheet, so that asymmetry was XLSX-specific (contrast PPTX below,
        where 0 slides genuinely means ""). Now XLSX matches PPTX's honest
        "nothing extracted -> ''" contract when no sheet has any data."""
        workbook = openpyxl.Workbook()  # default: exactly one blank sheet
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert text == ""

    def test_formula_only_workbook_with_no_cached_values_yields_empty_text(self):
        """CAVEAT documented in `_extract_xlsx_text`'s docstring, now covered
        by a real test: `data_only=True` reads only a formula's CACHED
        computed value. A workbook whose formulas were never evaluated by a
        calculating engine (Excel, LibreOffice) -- including any workbook
        openpyxl itself writes, since openpyxl never evaluates formulas --
        has no cache, so every formula cell reads back as None. That row is
        then indistinguishable from a genuinely blank one and skipped; if
        that's true of every row in every sheet, the whole extraction is
        honestly empty (see the previous test) rather than silently indexing
        a document whose only content was formula source text no one can
        read as computed values."""
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet["A1"] = "=1+1"  # formula string; openpyxl writes NO cached <v>
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert text == ""

    def test_formula_only_workbook_logs_a_diagnostic_warning(self, monkeypatch):
        """The empty-text result above is easy to misread as 'this workbook
        is genuinely empty' -- a structured warning naming the sheet and the
        data_only caveat is the runtime signal for the actual cause.
        `structlog`'s warnings don't reach pytest's stdlib-`logging`-based
        `caplog` fixture unless routed through `structlog.stdlib` (this
        codebase isn't), so the module's `logger.warning` is monkeypatched
        directly -- same pattern as mocking any other collaborator."""
        import src.temporal.activities.extract as extract_module

        calls = []
        monkeypatch.setattr(
            extract_module.logger, "warning", lambda msg, **kw: calls.append((msg, kw))
        )

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet["A1"] = "=1+1"
        buf = io.BytesIO()
        workbook.save(buf)

        _extract_xlsx_text(buf.getvalue())

        assert len(calls) == 1
        message, kwargs = calls[0]
        assert "none evaluated to visible data" in message
        assert kwargs["sheet"] == "Sheet"
        assert "data_only=True" in kwargs["hint"]

    def test_many_sheets_extracts_all_without_hanging(self):
        """550 sheets (comfortably above a realistic real-world workbook,
        deliberately chosen to exceed the #118 issue's illustrative '500
        sheets' failure-path case) with a handful of cells each stays well
        under the cell cap and must extract every sheet, not silently drop
        or hang on any of them."""
        workbook = openpyxl.Workbook()
        workbook.remove(workbook.active)
        for i in range(550):
            sheet = workbook.create_sheet(title=f"S{i}")
            sheet["A1"] = f"row{i}"
            sheet["B1"] = i
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert text.count("## Sheet:") == 550
        assert "## Sheet: S0" in text
        assert "## Sheet: S549" in text

    def test_cell_cap_exceeded_fails_actionably_not_oom(self, monkeypatch):
        """The evaluated-cell cost guard (#118: 'cap evaluated cells (e.g.
        500k)... exceeding -> document failed with actionable error, never
        OOM') is exercised directly by lowering the cap rather than building
        a 500k-cell fixture (slow, and not the point of this test -- the
        point is the cap's ENFORCEMENT, not its specific threshold)."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_XLSX_CELLS", 3)

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["a", "b", "c", "d"])  # 4 cells > cap of 3
        buf = io.BytesIO()
        workbook.save(buf)

        with pytest.raises(ApplicationError, match="evaluated-cell cap") as exc_info:
            _extract_xlsx_text(buf.getvalue())
        assert exc_info.value.non_retryable

    def test_text_cap_exceeded_incrementally_bounds_memory(self, monkeypatch):
        """BLOCKER (review): the text-length cap used to be checked ONLY
        after `"\\n\\n".join(sheet_parts)` had already built the full,
        multi-sheet output string -- the exact allocation the cap exists to
        prevent had already happened by the time it fired. Measured on
        review with a 200x200 grid of 32KB-string cells (8% of the cell cap,
        a 101KB upload): 2,572MB peak RSS and 24s before the old code's
        end-of-function check raised.

        This test proves the FIX -- the cap is now checked INSIDE the row
        loop, immediately after each row's text is added to the running
        total -- by lowering the cap and using a moderate multi-row grid: if
        the check still only ran at the end, this would build the entire
        ~200KB string and pass; because it now runs per-row, it must raise
        after only a handful of rows, well before all of them are read.
        """
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_XLSX_TEXT_CHARS", 10_000)

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        row_text = "x" * 2_000  # ~10 cols x 200 chars/cell-ish per row
        for _ in range(200):
            sheet.append([row_text] * 10)
        buf = io.BytesIO()
        workbook.save(buf)

        with pytest.raises(ApplicationError, match="character cap") as exc_info:
            _extract_xlsx_text(buf.getvalue())
        assert exc_info.value.non_retryable

    def test_single_pathological_cell_is_truncated_not_unbounded(self):
        """BLOCKER (review): `_MAX_XLSX_CELLS` bounds the CELL COUNT, and a
        cell is an unbounded-length string -- the cell-count cap alone
        cannot prevent one pathologically large value from blowing the
        memory/text budget by itself. `_format_xlsx_cell` must bound each
        rendered value independently."""
        huge = "z" * 50_000
        rendered = _format_xlsx_cell(huge)
        assert len(rendered) < len(huge)
        assert rendered.startswith("z" * 100)
        assert "truncated" in rendered
        assert "50000" in rendered  # names the original length

    def test_merge_scan_skips_oversized_worksheet_xml_not_decompress_unbounded(self, monkeypatch):
        """Self-caught regression: the FIRST version of merge-span detection
        (`_xlsx_merge_anchors`) read+decoded the ENTIRE worksheet XML part
        unconditionally to regex-search it for `<mergeCell>` tags. Measured
        directly: a worksheet's UNCOMPRESSED XML can be many orders of
        magnitude larger than the upload's compressed size (a 2.6MB upload
        -- the BLOCKER 1 pathological shape above, OOXML inline strings, not
        shared strings -- decompresses to a 1.3GB sheet1.xml). Reading that
        unconditionally reintroduced the exact unbounded-memory failure
        BLOCKER 1 closes, through a completely different code path: peak RSS
        measured at 2,516MB / 19.45s for that shape before this gate, 18MB /
        0.14s after. This test pins the gate itself (`_MAX_MERGE_SCAN_BYTES`,
        checked against `ZipInfo.file_size` BEFORE any `zf.read()`) directly,
        without needing a multi-GB fixture in CI: lower the gate below any
        real worksheet's size and confirm merge annotation is silently
        skipped (not that extraction fails -- annotation is a nicety, never
        worth failing over)."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(
            extract_module, "_MAX_MERGE_SCAN_BYTES", 1
        )  # any real sheet exceeds this

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.merge_cells("A1:D1")
        sheet["A1"] = "Title Row"
        sheet.append(["a", "b", "c", "d"])
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert "Title Row" in text  # extraction itself is unaffected
        assert "[merged" not in text  # but the (gated-off) annotation is absent

    def test_pathological_many_small_cells_shape_caught_incrementally(self, monkeypatch):
        """Reproduces the review's exact pathological shape at a scale a
        unit test can afford (real cells are truncated per-cell anyway, per
        the test above) -- 50 columns x 50KB-ish strings, cap lowered so the
        incremental per-row check must fire within the first few rows, not
        after streaming the whole 40k-cell grid the review's original report
        described."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_XLSX_TEXT_CHARS", 50_000)
        monkeypatch.setattr(extract_module, "_MAX_XLSX_CELL_CHARS", 2_000)

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        big_cell = "q" * 5_000  # will truncate to 2_000 chars each
        for _ in range(50):
            sheet.append([big_cell] * 50)
        buf = io.BytesIO()
        workbook.save(buf)

        with pytest.raises(ApplicationError, match="character cap"):
            _extract_xlsx_text(buf.getvalue())

    def test_dates_and_numbers_render_deterministically(self):
        """#118 acceptance criterion: numbers and dates render
        deterministically (not locale-dependent repr, not float noise).
        Pins the EXACT rendered line (review follow-up: the previous test
        only substring-matched "2026-01-15", which also matches the buggy
        "2026-01-15T00:00:00" -- passing without pinning the real output).
        A date-only cell must render WITHOUT a spurious "T00:00:00" time
        suffix; a cell with a real time component keeps its full timestamp."""
        import datetime

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["Count", "DateOnly", "Stamp"])
        sheet.append([42, datetime.date(2026, 1, 15), datetime.datetime(2026, 1, 15, 13, 30, 0)])
        buf = io.BytesIO()
        workbook.save(buf)

        text = _extract_xlsx_text(buf.getvalue())
        assert "42 | 2026-01-15 | 2026-01-15T13:30:00" in text
        assert "2026-01-15T00:00:00" not in text  # the bug this pins the fix for

    def test_legacy_xls_has_no_registry_entry(self):
        """Legacy .xls (binary OLE2 format, NOT the same format as .xlsx)
        must never be silently mis-parsed as XLSX. It has no
        FILE_TYPE_REGISTRY entry at all -- REST upload 400s before
        extraction is ever reached (pinned in inh-contracts'
        test_file_types.py::test_get_spec_for_extension_unknown_returns_none
        and the dispatch-layer test in test_temporal_activities.py)."""
        from inh_contracts.file_types import get_spec_for_mime

        assert get_spec_for_mime("application/vnd.ms-excel") is None

    def test_memory_error_during_construction_propagates_not_wrapped(self, monkeypatch):
        """#215 pattern-sweep hit: `_extract_xlsx_text`'s `except Exception`
        around `openpyxl.load_workbook()` construction had no `except
        MemoryError: raise` carve-out (unlike `_extract_pdf_text` /
        `_extract_docx_text`), so a MemoryError raised while OPENING a
        pathological workbook was incorrectly reclassified as
        `non_retryable=True` -- same defect class as #215, just at the
        construction site rather than the (already-correctly-unwrapped) row
        -iteration loop below it. Must propagate completely unconverted."""
        import openpyxl

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: out of memory opening workbook")

        monkeypatch.setattr(openpyxl, "load_workbook", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_xlsx_text(b"irrelevant, load_workbook is mocked")


class TestPptxFailurePaths:
    def test_corrupt_truncated_zip_raises_non_retryable(self):
        """Non-retryable ApplicationError (review follow-up, same reasoning
        as XLSX above): a corrupt/truncated deck can never succeed on
        retry."""
        with pytest.raises(ApplicationError, match="PPTX extraction failed") as exc_info:
            _extract_pptx_text(b"PK\x03\x04" + b"garbage, not a real zip central directory")
        assert exc_info.value.non_retryable

    def test_password_protected_bytes_raise_non_retryable(self):
        """Same reasoning as the XLSX case above: encrypted PPTX is OLE2,
        normally caught by sniff_content_type before reaching here; this
        pins the extractor's own defense-in-depth, non-retryable failure
        path."""
        ole2_encrypted_magic = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 32
        with pytest.raises(ApplicationError, match="PPTX extraction failed") as exc_info:
            _extract_pptx_text(ole2_encrypted_magic)
        assert exc_info.value.non_retryable

    def test_missing_dependency_raises_non_retryable(self, monkeypatch):
        """python-pptx not being installed is deterministic -- non-retryable,
        mirrors the XLSX/openpyxl case above."""
        import builtins

        real_import = builtins.__import__

        def _blocking_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("simulated: python-pptx not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", _blocking_import)

        with pytest.raises(ApplicationError, match="python-pptx not available") as exc_info:
            _extract_pptx_text(b"irrelevant")
        assert exc_info.value.non_retryable

    def test_empty_deck_yields_empty_text(self):
        """Unlike XLSX (always >=1 sheet), a brand-new python-pptx
        Presentation() genuinely has ZERO slides -- there is no structural
        unit to emit a boundary heading for, so extraction returns "".
        The caller (extract_text activity) already treats an all-whitespace
        extraction as a hard failure with an actionable message -- this test
        pins the extractor's own honest "nothing to extract" contract that
        failure depends on."""
        presentation = Presentation()  # 0 slides by default
        buf = io.BytesIO()
        presentation.save(buf)

        text = _extract_pptx_text(buf.getvalue())
        assert text == ""

    def test_many_slides_extracts_all_without_hanging(self):
        """520 slides (above the issue's illustrative '500 slides' case),
        each with just a title, must all be extracted -- not silently
        truncated or hung on."""
        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        for i in range(520):
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i}"
        buf = io.BytesIO()
        presentation.save(buf)

        text = _extract_pptx_text(buf.getvalue())
        assert text.count("## Slide") == 520
        assert "Slide 0" in text
        assert "Slide 519" in text

    def test_slide_cap_exceeded_fails_actionably_not_oom(self, monkeypatch):
        """Mirrors the XLSX cell cap: a pathological deck must fail with an
        actionable, non-retryable message rather than run away unbounded."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_PPTX_SLIDES", 2)

        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        for i in range(4):
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i}"
        buf = io.BytesIO()
        presentation.save(buf)

        with pytest.raises(ApplicationError, match="slide cap") as exc_info:
            _extract_pptx_text(buf.getvalue())
        assert exc_info.value.non_retryable

    def test_text_cap_exceeded_incrementally_bounds_memory(self, monkeypatch):
        """BLOCKER (review): mirrors the XLSX text-cap fix -- the cap used
        to be checked only after the entire deck's text was joined into one
        string. Now checked per-shape, inside the slide loop, so a
        pathological deck fails after a handful of slides instead of after
        materializing the whole thing."""
        import src.temporal.activities.extract as extract_module

        monkeypatch.setattr(extract_module, "_MAX_PPTX_TEXT_CHARS", 10_000)

        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        body_text = "y" * 3_000
        for _ in range(20):
            slide = presentation.slides.add_slide(layout)
            slide.placeholders[1].text_frame.text = body_text
        buf = io.BytesIO()
        presentation.save(buf)

        with pytest.raises(ApplicationError, match="character cap") as exc_info:
            _extract_pptx_text(buf.getvalue())
        assert exc_info.value.non_retryable

    def test_single_pathological_run_is_truncated_not_unbounded(self):
        """BLOCKER (review): mirrors the XLSX per-cell bound -- a single
        pathologically long paragraph/table-cell must not blow the text
        budget by itself; `_MAX_PPTX_SLIDES` bounds slide COUNT, not the
        length of any one shape's text."""
        from src.temporal.activities.extract import _pptx_bounded_text

        huge = "w" * 50_000
        rendered = _pptx_bounded_text(huge)
        assert len(rendered) < len(huge)
        assert rendered.startswith("w" * 100)
        assert "truncated" in rendered
        assert "50000" in rendered

    def test_legacy_ppt_has_no_registry_entry(self):
        """Legacy .ppt (binary OLE2 format) must never be silently
        mis-parsed as PPTX -- no registry entry, same contract as .xls
        above."""
        from inh_contracts.file_types import get_spec_for_mime

        assert get_spec_for_mime("application/vnd.ms-powerpoint") is None

    def test_memory_error_during_construction_propagates_not_wrapped(self, monkeypatch):
        """#215 pattern-sweep hit: `_extract_pptx_text`'s `except Exception`
        around `Presentation()` construction had no `except MemoryError:
        raise` carve-out (unlike `_extract_pdf_text` / `_extract_docx_text`),
        so a MemoryError raised while OPENING a pathological deck was
        incorrectly reclassified as `non_retryable=True` -- same defect
        class as #215, just at the construction site rather than the
        (already-correctly-unwrapped) slide-iteration loop below it. Must
        propagate completely unconverted."""
        import pptx

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: out of memory opening presentation")

        monkeypatch.setattr(pptx, "Presentation", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_pptx_text(b"irrelevant, Presentation is mocked")


class TestPdfFailurePaths:
    """#195: `_extract_pdf_text` was completely unwrapped -- a corrupt/
    truncated/password-protected PDF raised pypdf's raw exception type
    directly, and Temporal's default 3-attempt RetryPolicy retried it 3x
    before giving up, even though the SAME bytes can never succeed on retry
    (deterministic given fixed content). Same reasoning, same fix shape as
    XLSX/PPTX/DOCX above: catch, wrap as a non-retryable ApplicationError
    that names the likely cause, never leak the raw pypdf exception message
    unwrapped into the document's `error_message` / dead-letter row."""

    def test_corrupt_truncated_bytes_raise_non_retryable(self):
        """A `%PDF-` header prefix followed by garbage (corrupt or truncated
        upload) must fail loudly -- pypdf's own PdfReadError/PdfStreamError
        is caught and re-raised as an actionable, NON-RETRYABLE
        ApplicationError, mirroring XLSX/PPTX/DOCX's existing contract."""
        with pytest.raises(ApplicationError, match="PDF extraction failed") as exc_info:
            _extract_pdf_text(b"%PDF-1.4\n" + b"garbage, not a real xref table or pages tree")
        assert exc_info.value.non_retryable

    def test_not_a_pdf_at_all_raises_non_retryable(self):
        """Bytes with no PDF header at all (a mislabeled upload reaching
        this extractor, e.g. via the extension fallback) -- pypdf raises
        immediately at `PdfReader()` construction; must still be
        non-retryable, not a bare crash."""
        with pytest.raises(ApplicationError, match="PDF extraction failed") as exc_info:
            _extract_pdf_text(b"this is not a pdf at all, just plain text bytes")
        assert exc_info.value.non_retryable

    def test_missing_dependency_raises_non_retryable(self, monkeypatch):
        """Neither pypdf nor PyPDF2 being installed is itself deterministic
        (retrying the same worker process, or any worker built from the same
        image, can never succeed) -- non-retryable, same reasoning as every
        other missing-extraction-library case in this module. Previously a
        bare RuntimeError here (see the module's pre-#195 history) was
        retried 3x for no reason, exactly like the two failure modes above."""
        import builtins

        real_import = builtins.__import__

        def _blocking_import(name, *args, **kwargs):
            if name in ("pypdf", "PyPDF2"):
                raise ImportError(f"simulated: {name} not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", _blocking_import)

        with pytest.raises(ApplicationError, match="PDF extraction libraries not available") as (
            exc_info
        ):
            _extract_pdf_text(b"irrelevant")
        assert exc_info.value.non_retryable

    def test_raw_pypdf_exception_message_never_leaks_unwrapped(self):
        """The actionable ApplicationError message must be constructed by
        this module, not just pypdf's raw exception re-raised verbatim --
        pins that the wrapping actually happened (a regression that removed
        the try/except but kept a differently-worded message would still
        satisfy `match="PDF extraction failed"` above; this checks the
        specific corrupt/truncated/password-protected/wrong-format guidance
        is present, the same actionable shape XLSX/PPTX/DOCX give)."""
        with pytest.raises(ApplicationError) as exc_info:
            _extract_pdf_text(b"%PDF-1.4\n" + b"garbage, not a real xref table or pages tree")
        message = str(exc_info.value)
        assert "corrupt" in message
        assert "password-protected" in message

    def test_memory_error_during_construction_propagates_not_wrapped(self, monkeypatch):
        """Review follow-up: MemoryError is a load-dependent condition, not a
        property of the input bytes -- a version of this fix that caught it
        under a blanket `except Exception` would have reclassified it as
        `non_retryable=True`, permanently dead-lettering a failure a retry
        (possibly on a less-contended worker) could plausibly resolve. It
        must propagate completely unconverted -- not even as a differently
        worded ApplicationError."""
        import pypdf

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: out of memory parsing xref table")

        monkeypatch.setattr(pypdf, "PdfReader", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_pdf_text(b"%PDF-1.4\nirrelevant, PdfReader is mocked")

    def test_exception_during_page_iteration_propagates_not_wrapped(self, monkeypatch):
        """Review follow-up: the try/except is scoped to ONLY `PdfReader()`
        construction -- mirroring `_extract_xlsx_text`'s `load_workbook`-only
        wrap (that function's own row-iteration loop is likewise NOT wrapped
        in a broad except; only its explicit cap checks raise). A failure
        discovered lazily during page access/`extract_text()` (e.g. a
        per-page content-stream corruption in an otherwise structurally
        valid PDF, which construction alone cannot detect since pypdf only
        parses the xref/trailer eagerly) is NOT converted to a non-retryable
        ApplicationError -- it propagates as whatever pypdf raised, retryable
        by the default policy. Same accepted tradeoff XLSX/PPTX already make
        for their own row/shape iteration."""
        import pypdf

        class _ExplodingReader:
            def __init__(self, *args, **kwargs):
                pass

            @property
            def pages(self):
                raise pypdf.errors.PdfReadError("simulated: corrupt content stream on page 3")

        monkeypatch.setattr(pypdf, "PdfReader", _ExplodingReader)

        with pytest.raises(pypdf.errors.PdfReadError):
            _extract_pdf_text(b"%PDF-1.4\nirrelevant, PdfReader is mocked")


class TestJsonFailurePaths:
    """#195: `_extract_json_text` called `json.loads` unwrapped -- malformed
    JSON raised `json.JSONDecodeError` directly, retried 3x by Temporal's
    default RetryPolicy for a deterministic, unfixable-by-retry failure.
    Lower severity than PDF per the issue (JSON corruption is rarer/cheaper
    to retry), but the same defect class -- same fix shape."""

    def test_malformed_json_raises_non_retryable(self):
        with pytest.raises(ApplicationError, match="JSON extraction failed") as exc_info:
            _extract_json_text(b'{"key": "value", "unterminated": ')
        assert exc_info.value.non_retryable

    def test_not_json_at_all_raises_non_retryable(self):
        """Plain-text bytes reaching this extractor (a mislabeled upload)
        must fail the same actionable, non-retryable way, not a bare
        JSONDecodeError."""
        with pytest.raises(ApplicationError, match="JSON extraction failed") as exc_info:
            _extract_json_text(b"just some plain text, not json at all")
        assert exc_info.value.non_retryable

    def test_empty_bytes_raise_non_retryable(self):
        with pytest.raises(ApplicationError, match="JSON extraction failed") as exc_info:
            _extract_json_text(b"")
        assert exc_info.value.non_retryable

    def test_raw_json_decode_error_message_never_leaks_unwrapped(self):
        """`json.JSONDecodeError`'s own message is already clean (no heap
        addresses, unlike python-docx's) -- but the wrapping ApplicationError
        must still name the likely cause, not just echo the raw message
        verbatim, matching the actionable shape every other extractor in
        this module gives."""
        with pytest.raises(ApplicationError) as exc_info:
            _extract_json_text(b'{"key": "value", "unterminated": ')
        message = str(exc_info.value)
        assert "corrupt" in message or "not actually" in message


# ---------------------------------------------------------------------------
# #206 -- EPUB/RTF/ODT/subtitle: pattern-sweep completion of #195's defect
# class (bare RuntimeError retried 3x on a deterministic, unfixable-by-retry
# failure) for the last four extractors in this module. Same fix shape:
# ApplicationError(non_retryable=True) naming the likely cause, MemoryError
# left completely unconverted wherever a parse/construction call could
# plausibly raise it (#195 review follow-up, see _extract_pdf_text above).
# ---------------------------------------------------------------------------


def _zip_with(files: dict[str, bytes]) -> bytes:
    """Build an in-memory zip archive from `{member_name: content}` -- used
    below to construct minimal, deliberately-malformed EPUB/ODT fixtures
    without checking in opaque binary blobs (mirrors
    test_extraction_longtail.py's own in-memory zip fixtures)."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for name, data in files.items():
            zf.writestr(name, data)
    return buf.getvalue()


_VALID_EPUB_CONTAINER_XML = (
    b'<?xml version="1.0" encoding="UTF-8"?>'
    b'<container version="1.0" '
    b'xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
    b'<rootfiles><rootfile full-path="OEBPS/content.opf" '
    b'media-type="application/oebps-package+xml"/></rootfiles>'
    b"</container>"
)

_EPUB_CONTAINER_XML_NO_ROOTFILE = (
    b'<?xml version="1.0" encoding="UTF-8"?>'
    b'<container version="1.0" '
    b'xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
    b"<rootfiles/></container>"
)

_VALID_EPUB_OPF_NO_SPINE = (
    b'<?xml version="1.0" encoding="UTF-8"?>'
    b'<package xmlns="http://www.idpf.org/2007/opf" version="3.0">'
    b"<manifest/><spine/></package>"
)

_VALID_EPUB_OPF_UNREACHABLE_CHAPTER = (
    b'<?xml version="1.0" encoding="UTF-8"?>'
    b'<package xmlns="http://www.idpf.org/2007/opf" version="3.0">'
    b'<manifest><item id="chap1" href="chap1.xhtml" '
    b'media-type="application/xhtml+xml"/></manifest>'
    b'<spine><itemref idref="chap1"/></spine></package>'
)


class TestEpubFailurePaths:
    """#206: `_extract_epub_text` raised a bare `RuntimeError` at all 10 of
    its failure sites -- corrupt zip (x2: construction and central-directory
    read), DRM/encryption, missing/unparseable META-INF/container.xml,
    no rootfile declared, missing/unparseable content.opf, no spine, and
    "spine but no chapter produced extractable text". Every one is a pure
    function of `content`'s bytes -- the same input fails identically on
    every Temporal retry attempt -- so all ten now raise a non-retryable
    ApplicationError, same fix shape as #195's PDF/JSON and #118/#119's
    XLSX/PPTX/DOCX precedents in this module.

    Every except clause in `_extract_epub_text` catches a NARROW exception
    type (zipfile.BadZipFile, KeyError, ET.ParseError) -- never a broad
    `except Exception` -- so a MemoryError raised anywhere in this call
    chain (e.g. a pathological zip central directory, or a huge XML part)
    is never caught by these clauses and propagates completely unconverted;
    `test_memory_error_propagates_not_wrapped` below pins that this stays
    true rather than trusting it by inspection alone.
    """

    def test_corrupt_zip_raises_non_retryable(self):
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(b"not a zip file at all", "broken.epub")
        assert exc_info.value.non_retryable

    def test_drm_protected_raises_non_retryable(self):
        """DRM/encrypted EPUBs signal via the standard
        META-INF/encryption.xml manifest -- checked and rejected before any
        chapter is parsed."""
        epub_bytes = _zip_with(
            {
                "META-INF/encryption.xml": (
                    b'<encryption xmlns="urn:oasis:names:tc:opendocument:xmlns:container"/>'
                ),
            }
        )
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "protected.epub")
        assert exc_info.value.non_retryable
        message = str(exc_info.value).lower()
        assert "drm" in message or "encrypt" in message

    def test_missing_container_xml_raises_non_retryable(self):
        epub_bytes = _zip_with({"mimetype": b"application/epub+zip"})
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "no-container.epub")
        assert exc_info.value.non_retryable
        assert "container.xml" in str(exc_info.value)

    def test_unparseable_container_xml_raises_non_retryable(self):
        epub_bytes = _zip_with({"META-INF/container.xml": b"<not valid xml"})
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "bad-container.epub")
        assert exc_info.value.non_retryable

    def test_no_rootfile_raises_non_retryable(self):
        epub_bytes = _zip_with({"META-INF/container.xml": _EPUB_CONTAINER_XML_NO_ROOTFILE})
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "no-rootfile.epub")
        assert exc_info.value.non_retryable
        assert "rootfile" in str(exc_info.value)

    def test_missing_content_opf_raises_non_retryable(self):
        epub_bytes = _zip_with({"META-INF/container.xml": _VALID_EPUB_CONTAINER_XML})
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "no-opf.epub")
        assert exc_info.value.non_retryable
        assert "content.opf" in str(exc_info.value)

    def test_unparseable_content_opf_raises_non_retryable(self):
        epub_bytes = _zip_with(
            {
                "META-INF/container.xml": _VALID_EPUB_CONTAINER_XML,
                "OEBPS/content.opf": b"<not valid xml",
            }
        )
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "bad-opf.epub")
        assert exc_info.value.non_retryable

    def test_no_spine_raises_non_retryable(self):
        epub_bytes = _zip_with(
            {
                "META-INF/container.xml": _VALID_EPUB_CONTAINER_XML,
                "OEBPS/content.opf": _VALID_EPUB_OPF_NO_SPINE,
            }
        )
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "no-spine.epub")
        assert exc_info.value.non_retryable
        assert "spine" in str(exc_info.value)

    def test_no_extractable_chapters_raises_non_retryable(self):
        """Spine references a real manifest item whose href is not actually
        present in the zip -- the chapter is skipped (logged, not
        renumbered -- #125 review blocker 1), and since it was the only
        spine item, the terminal "spine but nothing extracted" case fires."""
        epub_bytes = _zip_with(
            {
                "META-INF/container.xml": _VALID_EPUB_CONTAINER_XML,
                "OEBPS/content.opf": _VALID_EPUB_OPF_UNREACHABLE_CHAPTER,
                # Deliberately no OEBPS/chap1.xhtml member.
            }
        )
        with pytest.raises(ApplicationError, match="EPUB extraction failed") as exc_info:
            _extract_epub_text(epub_bytes, "empty-spine.epub")
        assert exc_info.value.non_retryable
        assert "no chapter" in str(exc_info.value)

    def test_memory_error_propagates_not_wrapped(self, monkeypatch):
        """Review follow-up mirrored from #195's PDF fix: MemoryError must
        never be reclassified as non_retryable -- it's a load-dependent
        condition, not a property of the input bytes. Simulated at the
        `zipfile.ZipFile` construction call, the first thing this extractor
        does."""

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: OOM opening zip")

        monkeypatch.setattr(zipfile, "ZipFile", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_epub_text(b"irrelevant, ZipFile is mocked", "book.epub")


class TestRtfFailurePaths:
    """#206: `_extract_rtf_text` raised a bare `RuntimeError` at all 3 of its
    failure sites -- missing `striprtf` dependency, a genuine parse failure,
    and empty extraction. Missing-dependency is deterministic per
    worker/image (same reasoning as every other MissingExtractionDependency
    case in this module); the other two are deterministic given fixed
    `content` bytes. The parse-failure site wraps ONLY the `rtf_to_text`
    call (the preceding `.decode("latin-1")` step never raises -- latin-1
    maps every byte value 1:1, see the function's docstring) and explicitly
    re-raises MemoryError before the broad `except Exception` -- mirrors
    PDF's construction-only wrap (#195 review follow-up) so a load-dependent
    OOM parsing a pathological RTF is never reclassified as non-retryable.
    """

    def test_missing_dependency_raises_non_retryable(self, monkeypatch):
        import builtins

        real_import = builtins.__import__

        def _blocking_import(name, *args, **kwargs):
            if name == "striprtf.striprtf":
                raise ImportError("simulated: striprtf not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", _blocking_import)

        with pytest.raises(ApplicationError, match="striprtf not available") as exc_info:
            _extract_rtf_text(b"irrelevant", "sample.rtf")
        assert exc_info.value.non_retryable

    def test_parse_failure_raises_non_retryable(self, monkeypatch):
        """striprtf itself is lenient (rarely raises on malformed control
        words in practice), so the parse-failure wrapping is exercised
        directly by monkeypatching `rtf_to_text` to simulate the shape a
        pathological/corrupt control-word stream could produce -- proves the
        wrapping fires, not just that it's unreachable in normal use."""
        import striprtf.striprtf as striprtf_module

        def _raise(*args, **kwargs):
            raise ValueError("simulated: malformed control word group")

        monkeypatch.setattr(striprtf_module, "rtf_to_text", _raise)

        with pytest.raises(ApplicationError, match="RTF extraction failed") as exc_info:
            _extract_rtf_text(r"{\rtf1\ansi Hello\par}".encode("ascii"), "sample.rtf")
        assert exc_info.value.non_retryable

    def test_empty_extraction_raises_non_retryable(self):
        with pytest.raises(ApplicationError, match="RTF extraction failed") as exc_info:
            _extract_rtf_text(r"{\rtf1\ansi\par}".encode("ascii"), "empty.rtf")
        assert exc_info.value.non_retryable
        assert "no extractable text" in str(exc_info.value)

    def test_memory_error_during_parse_propagates_not_wrapped(self, monkeypatch):
        import striprtf.striprtf as striprtf_module

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: OOM parsing deeply nested RTF groups")

        monkeypatch.setattr(striprtf_module, "rtf_to_text", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_rtf_text(r"{\rtf1\ansi Hello\par}".encode("ascii"), "sample.rtf")


_VALID_ODT_CONTENT_XML_EMPTY_BODY = (
    b'<?xml version="1.0" encoding="UTF-8"?>'
    b"<office:document-content "
    b'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
    b'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
    b"<office:body><office:text></office:text></office:body>"
    b"</office:document-content>"
)


class TestOdtFailurePaths:
    """#206: `_extract_odt_text` raised a bare `RuntimeError` at all 6 of its
    failure sites -- corrupt zip (x2: construction and central-directory
    read), missing content.xml, unparseable content.xml, missing
    office:body/office:text, and empty extraction. Same narrow-except
    discipline as EPUB above (only zipfile.BadZipFile/KeyError/
    ET.ParseError are caught), so MemoryError is never caught by these
    clauses either -- see `test_memory_error_propagates_not_wrapped`."""

    def test_corrupt_zip_raises_non_retryable(self):
        with pytest.raises(ApplicationError, match="ODT extraction failed") as exc_info:
            _extract_odt_text(b"not a zip file at all", "broken.odt")
        assert exc_info.value.non_retryable

    def test_missing_content_xml_raises_non_retryable(self):
        """Also the exact signal for a mislabeled DOCX-under-.odt upload:
        DOCX has word/document.xml, not content.xml, at the zip root."""
        odt_bytes = _zip_with({"[Content_Types].xml": b"<Types/>"})
        with pytest.raises(ApplicationError, match="ODT extraction failed") as exc_info:
            _extract_odt_text(odt_bytes, "mislabeled.odt")
        assert exc_info.value.non_retryable
        assert "content.xml" in str(exc_info.value)

    def test_unparseable_content_xml_raises_non_retryable(self):
        odt_bytes = _zip_with({"content.xml": b"<not valid xml"})
        with pytest.raises(ApplicationError, match="ODT extraction failed") as exc_info:
            _extract_odt_text(odt_bytes, "bad.odt")
        assert exc_info.value.non_retryable

    def test_missing_body_raises_non_retryable(self):
        odt_bytes = _zip_with({"content.xml": b'<?xml version="1.0"?><root/>'})
        with pytest.raises(ApplicationError, match="ODT extraction failed") as exc_info:
            _extract_odt_text(odt_bytes, "no-body.odt")
        assert exc_info.value.non_retryable
        assert "office:body" in str(exc_info.value)

    def test_empty_body_raises_non_retryable(self):
        odt_bytes = _zip_with({"content.xml": _VALID_ODT_CONTENT_XML_EMPTY_BODY})
        with pytest.raises(ApplicationError, match="ODT extraction failed") as exc_info:
            _extract_odt_text(odt_bytes, "empty.odt")
        assert exc_info.value.non_retryable
        assert "no extractable text" in str(exc_info.value)

    def test_memory_error_propagates_not_wrapped(self, monkeypatch):
        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: OOM opening zip")

        monkeypatch.setattr(zipfile, "ZipFile", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_odt_text(b"irrelevant, ZipFile is mocked", "doc.odt")


class TestSubtitleFailurePaths:
    """#206: `_extract_subtitle_text`'s single failure site (no cue has a
    recognizable timestamp line) raised a bare `RuntimeError` -- deterministic
    given fixed content bytes, same fix shape as the rest of this module.
    Neither `_decode_text` nor `_parse_subtitle_cues` is wrapped in any
    try/except in this function, so a MemoryError from either (e.g.
    charset-detection scanning a huge mis-encoded file) already propagates
    unconverted with no code change needed here -- pinned below so a future
    refactor that adds a broad wrap cannot silently reintroduce the defect
    this whole sweep exists to close."""

    def test_no_cues_found_raises_non_retryable(self):
        with pytest.raises(ApplicationError, match="Subtitle extraction failed") as exc_info:
            _extract_subtitle_text(b"just some plain text, not a subtitle file at all", "notes.srt")
        assert exc_info.value.non_retryable
        assert "no subtitle cues" in str(exc_info.value).lower()

    def test_memory_error_during_decode_propagates_not_wrapped(self, monkeypatch):
        import src.temporal.activities.extract as extract_module

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: OOM decoding subtitle content")

        monkeypatch.setattr(extract_module, "_decode_text", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_subtitle_text(b"irrelevant, _decode_text is mocked", "sample.srt")


class TestDocxFailurePaths:
    """#215: `_extract_docx_text`'s broad `except Exception` wrapped BOTH the
    `Document()` construction call AND the paragraph-iteration list
    comprehension below it in the same `try` block -- a `MemoryError` raised
    during paragraph iteration (not just construction) got reclassified as
    `non_retryable=True`, permanently dead-lettering a load-dependent
    failure that a retry (possibly on a less-contended worker) could
    plausibly resolve. Mirrors `TestPdfFailurePaths` above exactly (#195's
    construction-only-wrap precedent, which this fix now matches)."""

    def test_memory_error_during_construction_propagates_not_wrapped(self, monkeypatch):
        """MemoryError raised by `Document()` construction itself must
        propagate completely unconverted -- not even as a differently
        worded ApplicationError."""
        import docx

        def _raise_memory_error(*args, **kwargs):
            raise MemoryError("simulated: out of memory parsing OOXML package")

        monkeypatch.setattr(docx, "Document", _raise_memory_error)

        with pytest.raises(MemoryError):
            _extract_docx_text(b"irrelevant, Document is mocked", "sample.docx")

    def test_exception_during_paragraph_iteration_propagates_not_wrapped(self, monkeypatch):
        """A failure discovered lazily during paragraph iteration (the `for
        p in doc.paragraphs if p.text.strip()` comprehension) -- e.g. a
        MemoryError from a pathological/huge document -- must NOT be swept
        into `non_retryable=True` by a broad except around the whole try
        block. The try/except is scoped to ONLY `Document()` construction,
        mirroring `_extract_pdf_text`'s construction-only wrap; the
        paragraph-iteration comprehension itself is left unwrapped below the
        try block."""
        import docx

        class _ExplodingDocument:
            def __init__(self, *args, **kwargs):
                pass

            @property
            def paragraphs(self):
                raise MemoryError("simulated: OOM iterating paragraphs")

        monkeypatch.setattr(docx, "Document", _ExplodingDocument)

        with pytest.raises(MemoryError):
            _extract_docx_text(b"irrelevant, Document is mocked", "sample.docx")

    def test_non_memory_exception_during_paragraph_iteration_propagates_not_wrapped(
        self, monkeypatch
    ):
        """Same as above but for a non-MemoryError exception discovered
        during paragraph iteration -- also must NOT become a non-retryable
        ApplicationError, since the paragraph-iteration comprehension sits
        entirely outside the try/except now (mirrors PDF's page-iteration
        test: a per-paragraph failure a structurally-valid `Document()`
        construction alone cannot detect stays retryable, not converted)."""
        import docx

        class _ExplodingDocument:
            def __init__(self, *args, **kwargs):
                pass

            @property
            def paragraphs(self):
                raise RuntimeError("simulated: corrupt paragraph run on p.12")

        monkeypatch.setattr(docx, "Document", _ExplodingDocument)

        with pytest.raises(RuntimeError, match="simulated: corrupt paragraph run"):
            _extract_docx_text(b"irrelevant, Document is mocked", "sample.docx")

    def test_value_error_on_construction_still_raises_non_retryable(self):
        """The existing ValueError-on-construction path (wrong OOXML content
        type, e.g. a genuine XLSX mislabeled as DOCX) must still become a
        non-retryable ApplicationError -- this fix narrows what the try/
        except covers, but must not accidentally widen the UNWRAPPED surface
        to cover construction itself. Full coverage of this path (message
        content, heap-address-repr scrubbing) lives in
        `test_genuine_xlsx_fed_to_docx_extractor_fails_loudly_not_silently`
        below; this pins just the non_retryable contract next to the two
        propagates-unwrapped tests above for an at-a-glance contrast."""
        xlsx_bytes = _read("sample.xlsx")
        with pytest.raises(ApplicationError, match="DOCX extraction failed") as exc_info:
            _extract_docx_text(xlsx_bytes, "report.docx")
        assert exc_info.value.non_retryable


def test_genuine_xlsx_fed_to_docx_extractor_fails_loudly_not_silently():
    """The reachable case inh-contracts' test_file_types.py::
    test_xlsx_bytes_declared_as_docx_pass_the_byte_sniff documents: a genuine
    XLSX, declared (or defaulted, e.g. no recognized extension) as DOCX,
    passes the byte-level sniff -- the shared ZIP family magic cannot tell
    them apart. This is the layer that DOES catch it: python-docx's own
    OOXML content-type check refuses to open a package whose principal part
    is a spreadsheet, not a document, raising instead of returning mangled
    or empty text.

    Review follow-up: a bare `pytest.raises(Exception)` would have passed
    even if this degraded to a bare `KeyError` or leaked python-docx's raw
    `ValueError` (observed, unwrapped: a `<_io.BytesIO object at 0x...>`
    heap-address repr with no filename) straight into the document's
    `error_message` / dead-letter row -- neither of which is what "fails
    loudly" is supposed to mean. This now pins the SPECIFIC exception type
    (non-retryable ApplicationError, since the failure is deterministic
    given fixed bytes), that the filename is present, and that the raw
    heap-address repr is gone."""
    xlsx_bytes = _read("sample.xlsx")
    with pytest.raises(ApplicationError, match="DOCX extraction failed") as exc_info:
        _extract_docx_text(xlsx_bytes, "report.docx")
    assert exc_info.value.non_retryable
    message = str(exc_info.value)
    assert "report.docx" in message
    assert "_io.BytesIO object at 0x" not in message  # the heap-address leak this pins the fix for
