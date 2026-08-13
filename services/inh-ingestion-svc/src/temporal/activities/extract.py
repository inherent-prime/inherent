"""Text extraction activity for converting document content to plain text.

Fetches file content directly from storage (instead of receiving bytes
via gRPC) and writes extracted text to the staging table.

Extraction dispatch (#117)
---------------------------
Which function handles a content type used to be an if/elif chain here,
duplicating the allow-list REST/MCP validation maintained independently in
``inh-public-api-svc``. Dispatch is now driven by the shared
``inh_contracts.FILE_TYPE_REGISTRY`` (the same registry REST/MCP validate
against): ``_resolve_extractor`` looks up the registry entry for a content
type and then the function wired for it in ``EXTRACTORS`` below. Two
failure modes are explicit and tested (see ``test_temporal_activities.py::
TestFileTypeRegistryDispatch``), never a silent lossy decode:

- No registry entry for the content type at all -> the document fails with
  a message naming the type and the supported set.
- A registry entry exists but its ``extractor`` key has no function in
  ``EXTRACTORS`` -- a wiring bug (a sibling format issue added a
  ``FileTypeSpec`` without its extractor) -- fails with a message that says
  so, instead of a bare ``KeyError`` crashing the Temporal worker.
"""

import datetime
import email
import io
import re
import zipfile
from collections.abc import Callable
from email import policy
from urllib.parse import unquote

# The three ET.fromstring() calls below parse XML pulled out of a
# customer-uploaded archive (EPUB container.xml / content.opf, ODT
# content.xml), so bandit flags the import as B405 and each call as B314.
# CPython >= 3.7.1 does not resolve external entities, which rules out classic
# XXE, but a hostile archive can still drive entity-expansion CPU burn -- so
# each call carries a `# nosec B314` pointing here, and #247 tracks moving all
# three to defusedxml (a runtime dependency add, out of scope for the CI
# change that first surfaced these).
from xml.etree import ElementTree as ET  # nosec B405 -- see above

import charset_normalizer
import structlog
from inh_contracts.file_types import all_mime_types, get_spec_for_upload
from temporalio import activity
from temporalio.exceptions import ApplicationError

from src.temporal.lineage import track_event
from src.temporal.models import ExtractTextInput, ExtractTextOutput

logger = structlog.get_logger(__name__)


@activity.defn
async def extract_text(input: ExtractTextInput) -> ExtractTextOutput:
    """Extract text from document content based on content type.

    Which formats are supported, and which function handles each, is defined
    once in ``inh_contracts.FILE_TYPE_REGISTRY`` (#117) -- see the module
    docstring above and ``docs/reference/file-types.md`` for the current
    list, rather than duplicated here where it would drift.

    The activity fetches file content from storage itself (avoiding the
    4MB gRPC limit) and writes extracted text to the staging table.

    Args:
        input: Contains storage refs, content type, filename, and workflow_run_id

    Returns:
        ExtractTextOutput with text_length (text itself is in staging)
    """
    async with track_event(
        workflow_run_id=input.workflow_run_id,
        document_id=input.document_id or "",
        workspace_id=input.workspace_id,
        event_type="text_extracted",
    ):
        return await _extract_text_inner(input)


async def _extract_text_inner(input: ExtractTextInput) -> ExtractTextOutput:
    """Inner implementation for text extraction (wrapped by lineage tracking)."""
    from src.temporal.shared_services import get_staging_service, get_storage_service

    # Fetch file content from storage
    storage_service = get_storage_service()

    if input.storage_backend == "local":
        content = storage_service.read_file(
            path=input.storage_path,
            backend="local",
            bucket=input.storage_bucket,
        )
    elif input.storage_backend == "gcs":
        content = storage_service.read_file(
            path=input.storage_path,
            backend="gcs",
            bucket=input.storage_bucket,
        )
    elif input.storage_backend == "s3":
        content = storage_service.read_file(
            path=input.storage_path,
            backend="s3",
            bucket=input.storage_bucket,
        )
    elif input.storage_backend == "azure":
        # #214: same gate as fetch.py's azure branch -- keep both in sync.
        # This activity runs AFTER fetch_document in the workflow, but must
        # not assume fetch_document's own gate already ran (activities are
        # independently retryable/replayable and must each be safe to call
        # on their own; see docs/developer/learnings.md on trusting an
        # earlier step's check instead of re-checking at the point of use).
        from src.config.settings import get_settings

        if not get_settings().allow_url_based_ingestion:
            raise RuntimeError(
                "Storage backend 'azure' (direct URL fetch) is disabled "
                "(ALLOW_URL_BASED_INGESTION is not set) -- see #214"
            )
        if input.storage_url:
            content = storage_service.read_file_from_url(input.storage_url)
        else:
            raise RuntimeError(f"Storage backend '{input.storage_backend}' requires storage_url")
    else:
        raise RuntimeError(f"Unknown storage backend: {input.storage_backend}")

    if content is None:
        raise RuntimeError("Failed to fetch document content from storage")

    content_type = input.content_type.lower()
    filename = input.original_filename.lower()

    logger.info(
        "Extracting text",
        content_type=content_type,
        filename=filename,
        content_size=len(content),
    )

    # Dispatch via the shared FILE_TYPE_REGISTRY (#117) -- see the module
    # docstring. `_resolve_extractor` raises a specific, actionable
    # RuntimeError for either "no registry entry" or "registry entry with no
    # wired extractor"; there is deliberately no catch-all decode-as-text
    # fallback anymore -- an unrecognized type must fail the document, never
    # silently produce garbled chunks. `original_filename` is passed through
    # so a generic/absent content type (e.g. "application/octet-stream")
    # persisted by the #122 upload-time extension fallback resolves the same
    # way here that it did at intake -- see `_resolve_extractor`.
    extractor = _resolve_extractor(content_type, input.original_filename)
    text = extractor(content, input.original_filename)

    # Run data quality checks on extracted text
    from src.services.quality import DataQualityService

    quality = DataQualityService()
    quality_results = quality.check_extracted_text(text, input.original_filename)
    quality.log_results(quality_results, document_id="extract:" + input.workflow_run_id)
    if quality.has_critical_failure(quality_results):
        raise RuntimeError(
            f"Text quality check failed for {input.original_filename}: empty extraction"
        )

    # Strip NUL (0x00) bytes before staging (issue #84). Postgres text/varchar
    # columns cannot store the NUL byte at all -- the driver raises before the
    # query reaches the server -- so an unsanitized value fails the staging
    # write permanently. Some documents (e.g. PDFs pypdf decodes imperfectly)
    # produce embedded NUL bytes. We strip *after* the quality check so its
    # `no_binary_content` diagnostic still sees the raw signal, and before the
    # empty-text guard so a text made up entirely of NUL bytes is caught below.
    if "\x00" in text:
        null_count = text.count("\x00")
        text = text.replace("\x00", "")
        logger.warning(
            "Stripped NUL bytes from extracted text before staging",
            filename=input.original_filename,
            null_bytes_removed=null_count,
        )

    if not text.strip():
        raise RuntimeError(
            f"Text extraction produced empty result for {filename} "
            f"(content_type={content_type}, size={len(content)} bytes)"
        )

    logger.info(
        "Text extracted successfully",
        content_type=content_type,
        text_length=len(text),
    )

    # Write extracted text to staging
    staging = get_staging_service()
    staging.write_text(input.workflow_run_id, text)

    return ExtractTextOutput(text_length=len(text))


def _decode_text(content: bytes) -> str:
    """Decode raw bytes to text without ever silently dropping bytes (#117).

    1. Try strict UTF-8 first: the overwhelming common case for uploaded
       text/markdown/csv/html content, and the only decoding that's provably
       correct when it succeeds -- it also handles content with embedded NUL
       bytes exactly right (needed for the #84 NUL-stripping step above),
       where encoding-detection heuristics on short/ambiguous input can
       misfire (see ``test_temporal_activities.py::TestDecodeText``).
    2. If the bytes are NOT valid UTF-8, use charset-normalizer to detect the
       actual encoding (Windows-1252, UTF-16, ...) and decode with it. This
       replaces the previous ``content.decode("utf-8", errors="ignore")``,
       which silently DELETED every byte that wasn't valid UTF-8 -- data loss
       with no signal, not even a log line.
    3. If detection can't confidently identify an encoding either, fall back
       to UTF-8 with replacement characters (a visible mojibake marker)
       rather than silently vanishing the byte.
    """
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        pass

    best = charset_normalizer.from_bytes(content).best()
    if best is not None:
        return str(best)
    return content.decode("utf-8", errors="replace")


def _extract_json_text(content: bytes) -> str:
    """Parse JSON and pretty-print it as extractable text.

    ``json.loads`` accepts bytes directly (auto-detecting UTF-8/16/32 via any
    BOM present, per the JSON spec), so this needs no separate decode step.

    Raises:
        ApplicationError (non_retryable=True): `content` is not valid JSON
            (malformed syntax, or not JSON at all -- e.g. a mislabeled
            upload reaching this extractor). Deterministic given fixed
            `content` bytes -- retrying cannot change the outcome (#195,
            same reasoning as the XLSX/PPTX/DOCX/PDF open failures in this
            module), so this is non-retryable rather than burning
            Temporal's default 3-attempt retry budget on a guaranteed-repeat
            failure. `json.JSONDecodeError`'s own message is already clean
            (no heap addresses, unlike python-docx's -- see
            `_extract_docx_text`), so it is safe to include directly.
    """
    import json

    try:
        data = json.loads(content)
    except json.JSONDecodeError as e:
        raise ApplicationError(
            f"JSON extraction failed: could not parse the document ({e}). "
            f"The file may be corrupt, truncated, or not actually JSON "
            f"despite its declared type.",
            type="JsonParseFailed",
            non_retryable=True,
        ) from e
    return json.dumps(data, indent=2)


# Extraction dispatch table (#117): one entry per FileTypeSpec.extractor key
# in inh_contracts.FILE_TYPE_REGISTRY. Adding a sibling format (#118 XLSX,
# #119 PPTX, ...) means adding ONE FileTypeSpec entry (services/inh-contracts)
# and ONE function + entry here -- `test_every_registry_extractor_key_is_wired`
# fails CI if the two ever disagree. Every extractor has the uniform
# ``(content: bytes, filename: str) -> str`` signature `_resolve_extractor`
# dispatches through, even where a given extractor ignores `filename` --
# lambdas adapt the helpers below that predate this table and only take
# `content`, so their own signatures/tests (and imports elsewhere) stay
# untouched. Lambdas also defer name lookup to call time, so this table can
# sit above the helper functions it references without a definition-order
# NameError at import.
EXTRACTORS: dict[str, Callable[[bytes, str], str]] = {
    "text_passthrough": lambda content, filename: _decode_text(content),
    "json_pretty": lambda content, filename: _extract_json_text(content),
    "html": lambda content, filename: _extract_html_text(content),
    "pdf": lambda content, filename: _extract_pdf_text(content),
    "docx": lambda content, filename: _extract_docx_text(content, filename),
    "xlsx": lambda content, filename: _extract_xlsx_text(content),
    "pptx": lambda content, filename: _extract_pptx_text(content),
    "image_ocr": lambda content, filename: _extract_image_text(content, filename),
    # Long-tail formats (#124/#125/#126).
    "eml": lambda content, filename: _extract_eml_text(content, filename),
    "epub": lambda content, filename: _extract_epub_text(content, filename),
    "rtf": lambda content, filename: _extract_rtf_text(content, filename),
    "odt": lambda content, filename: _extract_odt_text(content, filename),
    # #121: XML shares HTML's tag-stripping path -- see `_extract_xml_text`
    # for why `html.parser` (not an XML/DTD-aware parser) is the deliberate,
    # XXE-safe choice.
    "xml": lambda content, filename: _extract_xml_text(content),
    # #127: SRT and WebVTT share one cue parser (see `_extract_subtitle_text`).
    "srt": lambda content, filename: _extract_subtitle_text(content, filename),
    "vtt": lambda content, filename: _extract_subtitle_text(content, filename),
}


def _resolve_extractor(content_type: str, filename: str = "") -> Callable[[bytes, str], str]:
    """Resolve the extractor function for `content_type`, or fail loudly.

    Two distinct, explicit failure modes -- both DETERMINISTIC (the same
    `content_type` will fail identically on every retry, since neither is a
    transient dependency/network condition), so both raise a non-retryable
    ``ApplicationError`` instead of a bare ``RuntimeError``: Temporal fails
    the activity after the FIRST attempt rather than burning the workflow's
    full retry budget (multiple attempts with backoff) on a bug retrying
    cannot fix (#117 review item 13) -- cheaper and faster to reach the same
    terminal `failed` document status. Contrast with the storage read
    failures in `_extract_text_inner` (genuinely transient -- a retry can
    plausibly succeed once the dependency recovers), which deliberately keep
    the default retry policy. A missing extraction library is NOT one of
    those: every such case in this module (pypdf/PyPDF2, python-docx,
    openpyxl, python-pptx) is deterministic per worker/image -- retrying the
    same build can never install a package -- and is itself a non-retryable
    ``ApplicationError`` (``MissingExtractionDependency``), same reasoning as
    here (#195 closed the last gap, PDF's own missing-library case).

    1. No FILE_TYPE_REGISTRY entry for this content type at all -- an
       unsupported or unrecognized upload reaching extraction. `filename` is
       consulted as a fallback for a generic/absent content type (e.g.
       "application/octet-stream") via ``get_spec_for_upload`` (#122) -- the
       SAME resolution `inh-public-api-svc`'s intake validation already
       applied, so a document accepted at upload through that fallback does
       not permanently fail here for the same reason it would have been
       rejected there. Defaults to "" (no fallback) so every pre-#122 call
       site -- including this module's own non-generic content types --
       behaves exactly as before.
    2. A registry entry exists but its ``extractor`` key has no function
       wired into EXTRACTORS -- a wiring bug (a sibling format issue added a
       FileTypeSpec without its extractor), not a bad upload, but it must
       still fail the document with an actionable message instead of a bare
       KeyError crashing the Temporal worker.
    """
    spec = get_spec_for_upload(content_type, filename)
    if spec is None:
        raise ApplicationError(
            f"No extractor registered for content type '{content_type}'. "
            f"Supported types: {', '.join(all_mime_types())}",
            type="UnregisteredContentType",
            non_retryable=True,
        )

    extractor = EXTRACTORS.get(spec.extractor)
    if extractor is None:
        raise ApplicationError(
            f"Registry entry '{spec.key}' ({content_type}) names extractor "
            f"'{spec.extractor}', which has no function wired in EXTRACTORS. "
            f"This is a wiring bug: add EXTRACTORS['{spec.extractor}'] in "
            f"src/temporal/activities/extract.py.",
            type="ExtractorWiringBug",
            non_retryable=True,
        )
    return extractor


def _extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF content.

    Deterministic given fixed `content` bytes -- retrying cannot change the
    outcome, so both failure modes below raise a non-retryable
    ``ApplicationError`` instead of a bare exception (#195, same reasoning
    already applied to XLSX/PPTX/DOCX in this module -- see
    `_extract_xlsx_text`/`_extract_pptx_text`/`_extract_docx_text`). Before
    this fix, `pypdf.PdfReader(...)` was completely unwrapped: a corrupt/
    truncated/password-protected PDF raised pypdf's raw exception type
    (``PdfReadError``, ``PdfStreamError``, ...) directly, and Temporal's
    default 3-attempt RetryPolicy retried it 3x before giving up on an
    outcome already known at attempt 1 -- verified empirically that all
    three of those shapes fail at `PdfReader()` CONSTRUCTION (pypdf parses
    the xref table/trailer eagerly), not later during page access.

    The try/except is scoped to ONLY the `PdfReader()` construction call --
    matching `_extract_xlsx_text`'s `load_workbook()` precedent exactly
    (that function's row-iteration loop is likewise NOT wrapped in a broad
    except; only its explicit cap checks raise). The page-iteration loop
    below is deliberately left UNWRAPPED (review follow-up: an earlier
    version of this fix wrapped the whole loop, which would have turned a
    `MemoryError` from a large/pathological PDF into `non_retryable=True` --
    permanently dead-lettering a load-dependent failure a retry, possibly on
    a less-contended worker, could plausibly resolve, instead of the
    transient failure it actually is). The tradeoff this accepts: a
    corruption localized to one page's content stream (structurally valid
    xref/trailer, broken per-page data) that pypdf only discovers lazily
    during `.pages`/`extract_text()` retries 3x under the default policy
    rather than failing once -- same accepted tradeoff XLSX/PPTX already
    make for their own row/shape iteration.

    Raises:
        ApplicationError (non_retryable=True): neither pypdf nor PyPDF2 is
            installed (deterministic per worker/image -- retrying the same
            build can never succeed), or pypdf can't construct a `PdfReader`
            from `content` at all (corrupt/truncated bytes, password-
            protected file, or a different binary format entirely reaching
            this extractor despite its declared PDF type). One clear,
            actionable message either way -- never a bare pypdf exception
            type leaking into the document's ``error_message`` / dead-letter
            row.
    """
    try:
        import pypdf
    except ImportError:
        try:
            import PyPDF2 as pypdf  # type: ignore[no-redef]  # noqa: N813
        except ImportError:
            raise ApplicationError(
                "PDF extraction libraries not available (pypdf or PyPDF2)",
                type="MissingExtractionDependency",
                non_retryable=True,
            )

    try:
        reader = pypdf.PdfReader(io.BytesIO(content))
    except MemoryError:
        # A load-dependent condition, not a property of the bytes -- must
        # stay retryable (see the docstring above). pypdf's own read path
        # can allocate heavily while parsing a pathological xref/object
        # stream; this must never be reclassified as non_retryable.
        raise
    except Exception as e:
        # Covers corrupt/truncated PDFs (pypdf.errors.PdfReadError,
        # PdfStreamError, EmptyFileError, ...), password-protected files,
        # and any other "pypdf couldn't open this" failure -- one clear
        # message instead of a library-specific exception type leaking out,
        # same contract as XLSX/PPTX's own open-failure wrapping.
        raise ApplicationError(
            f"PDF extraction failed: could not read the document "
            f"({type(e).__name__}: {e}). The file may be corrupt, truncated, "
            f"password-protected, or not actually a PDF despite its "
            f"declared type.",
            type="PdfOpenFailed",
            non_retryable=True,
        ) from e

    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)


def _extract_docx_text(content: bytes, filename: str = "") -> str:
    """Extract text from DOCX content.

    Wrapped (review follow-up on #118/#119) so a mismatched OOXML sibling --
    see inh_contracts.file_types's docx entry comment: the shared ZIP magic
    across docx/xlsx/pptx means a mislabeled upload CAN reach this function
    with, say, genuine XLSX bytes -- fails with a clear, filename-bearing
    message instead of leaking python-docx's raw exception (observed
    verbatim: ``ValueError: file '<_io.BytesIO object at 0x7f...>' is not a
    Word file, content type is '...spreadsheetml.sheet.main+xml'``) -- a
    bare heap-address repr with no filename -- into the document's
    ``error_message`` and the dead-letter row. This is the extraction-stage
    safety net #118/#119's PR description names for the mislabeled-OOXML
    case; leaving it unwrapped while wrapping the two new extractors was a
    pattern-sweep miss caught on review.

    Deterministic given fixed `content` bytes -- retrying cannot change the
    outcome, so failures raise a non-retryable ``ApplicationError``, the same
    reasoning as the XLSX/PPTX open/cap failures below (and the existing
    ``_resolve_extractor`` "no extractor"/"wiring bug" cases).
    """
    try:
        from docx import Document
    except ImportError:
        raise ApplicationError(
            "python-docx not available for DOCX extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        )

    label = f" ({filename})" if filename else ""
    try:
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        # python-docx's own "wrong OOXML content type" ValueError embeds a
        # raw `<_io.BytesIO object at 0x...>` repr -- a heap address -- in
        # its message (observed verbatim: "file '<_io.BytesIO object at
        # 0x...>' is not a Word file, content type is '...'"). NEVER relay
        # `str(e)` verbatim here: extract just the useful "content type is
        # '...'" fragment for that specific, expected shape, and fall back
        # to the exception's TYPE NAME ONLY (never its message) for anything
        # else -- the safe default once one exception type from this same
        # library has already demonstrated it leaks an object repr.
        import re

        content_type_match = re.search(r"content type is '([^']+)'", str(e))
        detail = (
            f"wrong OOXML content type ({content_type_match.group(1)})"
            if content_type_match
            else type(e).__name__
        )
        raise ApplicationError(
            f"DOCX extraction failed{label}: could not read the document "
            f"({detail}). The file may be corrupt, truncated, "
            f"password-protected, or actually a different OOXML format "
            f"(e.g. XLSX or PPTX) despite its declared type.",
            type="DocxOpenFailed",
            non_retryable=True,
        ) from e


# Cost guards for XLSX extraction (#118 issue requirement: "cap evaluated
# cells (e.g. 500k) and emitted text length; exceeding -> document `failed`
# with actionable error, never OOM"). All three are checked INCREMENTALLY,
# INSIDE the row loop -- not after joining the whole workbook into one string
# -- so a pathological workbook fails fast with bounded peak memory instead
# of fully materializing before any guard can fire. The per-cell bound
# matters as much as the aggregate cap: `_MAX_XLSX_CELLS` counts CELLS, and a
# cell is an unbounded-length string -- a grid at 8% of the cell cap with
# 32KB strings in every cell reached 2.5GB peak RSS before the aggregate
# checks (measured on review) because nothing bounded any single value.
_MAX_XLSX_CELLS = 500_000
_MAX_XLSX_CELL_CHARS = 10_000
_MAX_XLSX_TEXT_CHARS = 5_000_000

# Cheap chunking insurance for #129 (which owns the real per-format chunker):
# a 10k-row sheet flattens to one string where only the FIRST chunk a
# downstream fixed-size chunker produces carries the sheet name and header
# row -- re-emitting both every N data rows means most chunks still carry
# that context even before #129 lands.
_XLSX_HEADER_REPEAT_ROWS = 50

# Mirrors the XLSX guards above for PPTX: a slide-count ceiling generous
# enough that a genuinely large deck (the #119 issue's illustrative "500
# slides" case) still extracts in full, a per-run/per-cell character bound
# so one pathological paragraph or table cell can't blow the budget alone,
# and a text-length cap checked incrementally per slide (not after joining
# the whole deck) -- same reasoning as XLSX above.
_MAX_PPTX_SLIDES = 5_000
_MAX_PPTX_RUN_CHARS = 10_000
_MAX_PPTX_TEXT_CHARS = 5_000_000


def _format_xlsx_cell(value: object) -> str:
    """Render one cell's value deterministically (#118 acceptance criterion:
    "Numbers and dates render deterministically; formula cells render
    computed values"), bounded so one pathological value cannot blow the
    text-length cost guard alone.

    ``openpyxl`` (opened with ``data_only=True``, see `_extract_xlsx_text`)
    already resolves formula cells to their last-computed value before this
    function ever sees them, so there is no formula-vs-literal branch here --
    every value arriving is already the value to render.

    Dates: openpyxl has no separate "date-only" Python type -- a cell
    formatted in Excel as a pure date (no time component) still comes back
    as a `datetime.datetime` with the time fixed at midnight, not a
    `datetime.date`. Rendering every `datetime.datetime` with a full
    ISO-8601 timestamp would put a noisy, meaningless "T00:00:00" suffix on
    every date column. Instead: a `datetime.datetime` at exactly midnight
    renders as a bare date; anything with a real time component keeps the
    full timestamp. Accepted tradeoff, stated plainly: a genuine timestamp
    that happens to land exactly on midnight also renders as date-only --
    openpyxl gives no other signal at this API layer to distinguish the two
    without dropping ``read_only=True`` for a second, non-read-only parse
    (the per-cell number-format lookup `.iter_rows(values_only=True)`
    intentionally forgoes), which would undo the memory-safety that buys us.
    """
    if value is None:
        return ""
    if isinstance(value, datetime.datetime):
        rendered = (
            value.date().isoformat()
            if value.time() == datetime.time(0, 0, 0, 0)
            else value.isoformat()
        )
    elif isinstance(value, datetime.date):
        rendered = value.isoformat()
    else:
        rendered = str(value)

    if len(rendered) > _MAX_XLSX_CELL_CHARS:
        # One pathological cell must not blow the text-length budget alone
        # -- truncate with a visible marker (never a silent partial value).
        original_len = len(rendered)
        rendered = rendered[:_MAX_XLSX_CELL_CHARS] + f"...[truncated, {original_len} chars]"
    return rendered


# Uncompressed-size gate for `_xlsx_merge_anchors` below. A worksheet's XML
# can be many orders of magnitude larger UNCOMPRESSED than the upload's
# on-disk (compressed) size -- verified on review: a 2.6MB uploaded .xlsx
# (pathological, from the BLOCKER 1 fix above: a 200x200 grid of 32KB-string
# cells written as OOXML inline strings, not shared strings) decompresses to
# a 1.3GB `sheet1.xml`. `zipfile.ZipFile.read()` has no size limit of its
# own -- reading that part unconditionally would decompress the whole 1.3GB
# into memory just to regex-search it for `<mergeCell>` tags, silently
# reintroducing the exact unbounded-memory failure mode BLOCKER 1 exists to
# close, through a completely different code path. Any worksheet whose
# uncompressed XML exceeds this cap skips merge-span annotation entirely
# (falls back to `{}`, i.e. "no marker" -- a presentation nicety, never
# worth this risk) rather than reading it. In practice this never matters
# for a legitimate large sheet either: a sheet whose iter_rows() reaches
# `_MAX_XLSX_TEXT_CHARS` (which a sheet of this raw size will hit almost
# immediately) fails via that cap regardless of whether this function ran.
_MAX_MERGE_SCAN_BYTES = 5_000_000


def _xlsx_merge_anchors(content: bytes, worksheet_path: str | None) -> dict[str, str]:
    """Best-effort ``{anchor_coordinate: "A1:D1"}`` map of `worksheet_path`'s
    merged cell ranges, so a merged cell's rendered value can carry a
    ``[merged A1:D1]`` marker instead of flattening to a value cell followed
    by silently blank cells with nothing distinguishing "merged" from
    "genuinely empty".

    Read directly from the sheet's raw XML via a targeted regex, NOT via
    openpyxl's `Worksheet.merged_cells` API -- `ReadOnlyWorksheet` (what
    ``read_only=True`` gives us) has no such attribute at all (verified:
    accessing it raises `AttributeError`).

    Gated by `_MAX_MERGE_SCAN_BYTES` on the part's UNCOMPRESSED size (see
    that constant's comment) before ever calling `zf.read()` -- checking
    `ZipInfo.file_size` costs nothing (it's zip central-directory metadata,
    already in hand once the archive is opened) and is what makes this safe
    to call unconditionally, once per sheet, regardless of the upload's
    contents.

    Returns ``{}`` (no annotation, extraction proceeds unaffected) if the
    part is too large to scan safely, or if anything else about this lookup
    fails -- merge-span annotation is a presentation nicety, never worth
    failing (or slowing, or ballooning the memory of) an otherwise-successful
    extraction over.
    """
    if not worksheet_path:
        return {}
    try:
        import re
        import zipfile

        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            info = zf.getinfo(worksheet_path)
            if info.file_size > _MAX_MERGE_SCAN_BYTES:
                return {}
            xml = zf.read(worksheet_path).decode("utf-8", errors="ignore")
        anchors: dict[str, str] = {}
        for match in re.finditer(r'<mergeCell\s+ref="([A-Z]+\d+):[A-Z]+\d+"', xml):
            full_ref = match.group(0)
            ref = re.search(r'ref="([^"]+)"', full_ref)
            if ref is None:
                continue
            full_range = ref.group(1)
            anchors[match.group(1)] = full_range
        return anchors
    except Exception:
        return {}


def _extract_xlsx_text(content: bytes) -> str:
    """Extract text from XLSX content with row-aware, sheet-boundary
    serialization (#118).

    Per sheet, emits ``## Sheet: <name>`` then one pipe-delimited line per
    non-empty row, cells in column order -- so an agent reading the
    flattened text can still tell which value sat in which column (row-aware
    serialization) and which sheet a row came from (sheet boundaries). This
    property holds for the raw extracted STRING; it does NOT by itself
    survive a downstream fixed-size chunker splitting that string apart
    (measured on review: a 10k-row sheet produces ~669 chunks, of which
    exactly one carries the header row and exactly one carries the "## Sheet:"
    line) -- `_XLSX_HEADER_REPEAT_ROWS` below re-emits both periodically as
    cheap insurance ahead of #129's real per-format chunker.

    ``data_only=True`` reads each formula cell's last-COMPUTED value (the
    cached result Excel/LibreOffice stores when it saves the file) rather
    than the formula source text -- exactly the "computed values only, no
    formula source" contract in the #118 issue. CAVEAT: a workbook whose
    formulas were never cached by a calculating engine (Excel, LibreOffice --
    including any workbook openpyxl itself wrote, since openpyxl does not
    evaluate formulas) reads every formula cell as `None` -- a formula-only
    row is then indistinguishable from a genuinely blank one and is skipped
    (see the blank-row-skip comment below); a sheet logged as "had rows but
    none evaluated to data" is the runtime signal for exactly this case.
    ``read_only=True`` streams rows instead of loading the whole workbook
    into memory, which is what makes the incremental caps below actually
    protective instead of cosmetic.

    Raises:
        ApplicationError (non_retryable=True): openpyxl can't open `content`
            at all (corrupt/truncated zip, password-protected/OLE2 file, a
            legitimately different binary format sharing the OOXML zip
            signature -- see inh_contracts.file_types's docx entry comment),
            the evaluated-cell cap is exceeded, or the running character
            count exceeds the text cap. Deterministic given fixed bytes --
            retrying cannot change the outcome, so these are non-retryable
            rather than burning Temporal's default retry budget on a
            guaranteed-repeat failure. Every message is clear and actionable
            -- never a bare zipfile/openpyxl exception surfacing to the
            caller, and never a silent partial result.
    """
    try:
        import openpyxl
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise ApplicationError(
            "openpyxl not available for XLSX extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        )

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    except Exception as e:
        # Covers corrupt/truncated zips (zipfile.BadZipFile), password-
        # protected files (OLE2/CFBF container -- not a zip at all), and any
        # other "openpyxl couldn't make sense of this" failure -- one clear
        # message instead of a library-specific exception type leaking out.
        raise ApplicationError(
            f"XLSX extraction failed: could not open workbook ({type(e).__name__}: {e}). "
            f"The file may be corrupt, truncated, password-protected, or not "
            f"actually an XLSX despite its declared type.",
            type="XlsxOpenFailed",
            non_retryable=True,
        ) from e

    sheet_parts: list[str] = []
    total_cells = 0
    total_chars = 0
    any_data_emitted = False
    try:
        for sheet in workbook.worksheets:
            heading = f"## Sheet: {sheet.title}"
            sheet_lines = [heading]
            total_chars += len(heading)
            merge_anchors = _xlsx_merge_anchors(content, getattr(sheet, "_worksheet_path", None))

            header_line: str | None = None
            rows_since_heading = 0
            rows_with_data = 0
            row_number = sheet.min_row or 1
            min_column = sheet.min_column or 1

            for row in sheet.iter_rows(values_only=True):
                total_cells += len(row)
                if total_cells > _MAX_XLSX_CELLS:
                    raise ApplicationError(
                        f"XLSX extraction failed: evaluated-cell cap "
                        f"({_MAX_XLSX_CELLS}) exceeded while reading sheet "
                        f"'{sheet.title}'. Split the workbook into smaller "
                        f"files and re-upload.",
                        type="XlsxCellCapExceeded",
                        non_retryable=True,
                    )
                # Skip fully blank rows (read-only mode yields a None-filled
                # row for a genuinely blank row, and -- see the CAVEAT above
                # -- for a row of uncached formula cells too, which looks
                # identical at this layer).
                if all(cell is None for cell in row):
                    row_number += 1
                    continue

                cells = []
                for col_offset, cell in enumerate(row):
                    rendered = _format_xlsx_cell(cell)
                    coord = f"{get_column_letter(min_column + col_offset)}{row_number}"
                    if coord in merge_anchors:
                        rendered = f"{rendered} [merged {merge_anchors[coord]}]"
                    cells.append(rendered)
                line = " | ".join(cells)

                if header_line is None:
                    header_line = line
                elif rows_since_heading >= _XLSX_HEADER_REPEAT_ROWS:
                    repeat_heading = f"## Sheet: {sheet.title} (continued)"
                    sheet_lines.append(repeat_heading)
                    sheet_lines.append(header_line)
                    total_chars += len(repeat_heading) + len(header_line)
                    rows_since_heading = 0

                sheet_lines.append(line)
                total_chars += len(line)
                rows_since_heading += 1
                rows_with_data += 1
                any_data_emitted = True
                row_number += 1

                # Checked INSIDE the row loop, right after the line that
                # pushed it over -- not after joining the whole workbook --
                # so peak memory is bounded to a small, fixed multiple of the
                # cap (a few oversized rows), never the full pathological
                # input's worth.
                if total_chars > _MAX_XLSX_TEXT_CHARS:
                    raise ApplicationError(
                        f"XLSX extraction failed: extracted text exceeds the "
                        f"{_MAX_XLSX_TEXT_CHARS}-character cap (hit while "
                        f"reading sheet '{sheet.title}'). Split the workbook "
                        f"into smaller files and re-upload.",
                        type="XlsxTextCapExceeded",
                        non_retryable=True,
                    )

            if rows_with_data == 0 and (sheet.max_row or 0) > 0:
                logger.warning(
                    "XLSX sheet had rows but none evaluated to visible data",
                    sheet=sheet.title,
                    max_row=sheet.max_row,
                    hint=(
                        "If this sheet contains formulas, data_only=True reads "
                        "only their cached computed value -- a workbook never "
                        "opened/saved by a calculating engine (Excel, "
                        "LibreOffice) has no cache, so formula cells read as "
                        "empty. Re-save the file in a spreadsheet application "
                        "before re-uploading."
                    ),
                )

            sheet_parts.append("\n".join(sheet_lines))
    finally:
        # read_only workbooks hold an open zip/file handle until closed --
        # always release it, success or failure.
        workbook.close()

    if not any_data_emitted:
        # No sheet in the entire workbook produced a real data row -- return
        # an honestly empty extraction (not "## Sheet: Sheet" masquerading as
        # content) so the caller's existing empty-extraction guard fails the
        # document instead of silently indexing a content-free one. Mirrors
        # PPTX's "0 slides -> ''" contract below -- XLSX's asymmetry (every
        # workbook has >=1 sheet, unlike a deck's 0-slides case) previously
        # meant an empty/uncached-formula workbook cleared that guard by
        # accident via its sheet-heading text alone.
        return ""

    return "\n\n".join(sheet_parts)


def _pptx_slide_title(slide: object) -> str | None:
    """Best-effort slide title, or None if this slide has no title
    placeholder (a valid, common case -- e.g. a section-divider or
    image-only slide)."""
    shapes = getattr(slide, "shapes", None)
    title_shape = getattr(shapes, "title", None) if shapes is not None else None
    if title_shape is None:
        return None
    text = (title_shape.text or "").strip()
    return text or None


def _pptx_slide_notes(slide: object) -> str | None:
    """Speaker notes text for `slide`, or None if it has no notes slide, or
    the notes slide has no non-whitespace text."""
    # getattr rather than attribute access, matching _pptx_slide_title above:
    # `slide` is typed `object` because python-pptx ships no usable stubs.
    if not getattr(slide, "has_notes_slide", False):
        return None
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return None
    notes_text = (notes_slide.notes_text_frame.text or "").strip()
    return notes_text or None


def _pptx_bounded_text(text: str) -> str:
    """Bound a single paragraph/table-cell's rendered text so one
    pathological shape cannot blow the text-length budget alone -- mirrors
    `_format_xlsx_cell`'s per-cell bound, same reasoning (see the
    module-level cost-guard comment above `_MAX_XLSX_CELLS`)."""
    if len(text) > _MAX_PPTX_RUN_CHARS:
        original_len = len(text)
        return text[:_MAX_PPTX_RUN_CHARS] + f"...[truncated, {original_len} chars]"
    return text


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX content with slide-boundary sections (#119).

    Per slide, emits ``## Slide <n>: <title>`` (or ``## Slide <n>`` when the
    slide has no title placeholder), then every text-frame shape's text in
    shape order (reading order as authored), then any table shape's rows
    pipe-delimited (same row-aware convention as XLSX, #118), then speaker
    notes under a ``Notes:`` line -- so a query matching only speaker-notes
    text still lands in the same chunk as its slide's visible content once
    #129's chunker splits on these section boundaries. Embedded images are
    deliberately excluded in v1 (#119: "no OCR to manage costs" -- consistent
    with PNG's OCR being an explicit opt-in optional extra elsewhere in this
    module, not a default-on cost for every upload).

    Raises:
        ApplicationError (non_retryable=True): python-pptx can't open
            `content` at all (corrupt/truncated zip, password-protected/OLE2
            file, a different binary format sharing the OOXML zip
            signature), the slide-count cap is exceeded, or the running
            character count exceeds the text cap. Same "deterministic ->
            non-retryable, clear, actionable, never silent" contract as
            XLSX above.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ApplicationError(
            "python-pptx not available for PPTX extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        )

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as e:
        # Covers corrupt/truncated zips, password-protected (OLE2/CFBF)
        # files, and any other "python-pptx couldn't open this" failure.
        raise ApplicationError(
            f"PPTX extraction failed: could not open presentation "
            f"({type(e).__name__}: {e}). The file may be corrupt, truncated, "
            f"password-protected, or not actually a PPTX despite its "
            f"declared type.",
            type="PptxOpenFailed",
            non_retryable=True,
        ) from e

    slide_parts: list[str] = []
    total_chars = 0
    for index, slide in enumerate(presentation.slides, start=1):
        if index > _MAX_PPTX_SLIDES:
            raise ApplicationError(
                f"PPTX extraction failed: slide cap ({_MAX_PPTX_SLIDES}) "
                f"exceeded. Split the deck into smaller files and re-upload.",
                type="PptxSlideCapExceeded",
                non_retryable=True,
            )

        title = _pptx_slide_title(slide)
        heading = f"## Slide {index}: {title}" if title else f"## Slide {index}"
        slide_lines = [heading]
        total_chars += len(heading)

        for shape in slide.shapes:
            if shape.has_text_frame:
                # Title text is already in the heading above -- skip it here
                # so it isn't duplicated in the body.
                if shape == slide.shapes.title:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = _pptx_bounded_text("".join(run.text for run in paragraph.runs))
                    if paragraph_text.strip():
                        slide_lines.append(paragraph_text)
                        total_chars += len(paragraph_text)
            elif shape.has_table:
                for row in shape.table.rows:
                    line = " | ".join(_pptx_bounded_text(cell.text) for cell in row.cells)
                    slide_lines.append(line)
                    total_chars += len(line)

            # Checked INSIDE the shape loop, right after the text that pushed
            # it over -- not after joining the whole deck -- so peak memory
            # is bounded, same reasoning as XLSX's per-row check above.
            if total_chars > _MAX_PPTX_TEXT_CHARS:
                raise ApplicationError(
                    f"PPTX extraction failed: extracted text exceeds the "
                    f"{_MAX_PPTX_TEXT_CHARS}-character cap (hit on slide "
                    f"{index}). Split the deck into smaller files and "
                    f"re-upload.",
                    type="PptxTextCapExceeded",
                    non_retryable=True,
                )

        notes = _pptx_slide_notes(slide)
        if notes:
            notes = _pptx_bounded_text(notes)
            slide_lines.append("Notes:")
            slide_lines.append(notes)
            total_chars += len(notes) + len("Notes:")
            if total_chars > _MAX_PPTX_TEXT_CHARS:
                raise ApplicationError(
                    f"PPTX extraction failed: extracted text exceeds the "
                    f"{_MAX_PPTX_TEXT_CHARS}-character cap (hit on slide "
                    f"{index}'s notes). Split the deck into smaller files "
                    f"and re-upload.",
                    type="PptxTextCapExceeded",
                    non_retryable=True,
                )

        slide_parts.append("\n".join(slide_lines))

    return "\n\n".join(slide_parts)


def _extract_image_text(content: bytes, original_filename: str) -> str:
    """Extract text from a PNG image via Tesseract OCR with graceful fallback.

    OCR is an optional capability (requires the ``ocr`` extra plus the
    ``tesseract`` system binary). When OCR is unavailable -- the libraries
    are not installed, the tesseract binary is missing, or the image simply
    contains no readable text -- this returns a minimal placeholder instead
    of raising. The placeholder keeps the document flowing through the
    pipeline (0 useful chunks, but not a hard failure) so a missing OCR
    install never crashes ingestion.

    Args:
        content: Raw PNG bytes.
        original_filename: Original filename, used in the fallback placeholder.

    Returns:
        OCR-extracted text, or a placeholder string when OCR yields nothing
        or is unavailable.
    """
    placeholder = f"[image: {original_filename}, no text extracted]"

    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        logger.warning(
            "OCR libraries not available (install the 'ocr' extra: pytesseract, pillow); "
            "returning placeholder for image",
            filename=original_filename,
        )
        return placeholder

    try:
        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image)
    except pytesseract.TesseractNotFoundError:
        logger.warning(
            "Tesseract binary not found; install the 'tesseract-ocr' system package. "
            "Returning placeholder for image",
            filename=original_filename,
        )
        return placeholder
    except Exception as e:
        logger.warning(
            "OCR failed for image; returning placeholder",
            filename=original_filename,
            error=str(e),
        )
        return placeholder

    if not text.strip():
        logger.warning(
            "OCR produced no text for image; returning placeholder",
            filename=original_filename,
        )
        return placeholder

    return text


def _extract_html_text(content: bytes) -> str:
    """Extract text from HTML content.

    Falls back to a raw text decode if BeautifulSoup is not available (bs4 is
    a core, non-optional dependency of this service, so this branch is
    defense-in-depth rather than an expected runtime path). Uses
    `_decode_text` (#117) rather than `errors="ignore"` so even this fallback
    never silently drops bytes.
    """
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")
        for element in soup(["script", "style"]):
            element.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        logger.warning("beautifulsoup4 not available, falling back to raw decode")
        return _decode_text(content)


# ---------------------------------------------------------------------------
# #124 -- EML (RFC 822 email)
# ---------------------------------------------------------------------------


def _eml_attachment_label(part: email.message.EmailMessage) -> str:
    """A human-readable label for an EML attachment part, used only to name
    what was elided -- never its content (#124: "silently including their
    filenames as if they were content is not [defensible]" only applies to
    dumping them into the body; a clearly labeled filename listing is fine).
    """
    filename = part.get_filename()
    if filename:
        return filename
    if part.get_content_type() == "message/rfc822":
        # A forwarded/embedded email with no explicit filename -- label it
        # with its own Subject so it's still identifiable. Deliberately does
        # NOT look any deeper than this one header: #124 scopes nested
        # message/rfc822 parts to "first level only", and `iter_attachments`
        # below already only walks the OUTER message's immediate children,
        # so this function never gets asked about a grandchild part.
        try:
            nested = part.get_payload(0)
            # get_payload(0) is typed to possibly return a raw `str` (the
            # non-multipart-message case) as well as a Message -- only a
            # Message has headers to read `.get("Subject")` from.
            subject = nested.get("Subject") if isinstance(nested, email.message.Message) else None
        except (IndexError, AttributeError):
            subject = None
        return f"(embedded message: {subject})" if subject else "(embedded message)"
    return "(unnamed attachment)"


def _extract_eml_text(content: bytes, filename: str) -> str:
    """Extract text from an RFC 822 (.eml) email message (#124).

    An email is a tree, not a document -- three deliberate decisions:

    1. Headers (From/To/Cc/Date/Subject) are what make an email citable and
       searchable for an AI agent reader, so they are ALWAYS emitted first,
       never dropped -- even when there is no body at all.
    2. Body: prefer the text/plain part; fall back to text/html run through
       the EXISTING `_extract_html_text` (no second bespoke HTML parser).
       `email.policy.default`'s `EmailMessage.get_content()` already decodes
       quoted-printable/base64 Content-Transfer-Encoding and any declared
       charset for us -- no manual decode step needed here.
    3. Attachments (v1): NOT extracted. Their filenames (and a count) are
       recorded in a clearly labeled, separate section so an agent knows
       content was elided -- this is explicitly NOT the same as silently
       including filenames as if they were body content. Nested
       message/rfc822 parts are inspected one level only: `iter_attachments`
       walks the OUTER message's immediate children only, so an email
       forwarded as an attachment is listed by name but never recursed into
       for its own body/attachments.

    Raises nothing on its own for a missing body or missing headers -- an
    email with truly no headers, body, or attachments extracts to an empty
    string, and the caller's existing empty-text guard
    (`_extract_text_inner`) is what turns that into the actual document
    failure. This keeps exactly one place responsible for "empty extraction
    is a hard failure" instead of duplicating that check here.
    """
    msg = email.message_from_bytes(content, policy=policy.default)

    header_lines = [
        f"{header}: {value}"
        for header in ("From", "To", "Cc", "Date", "Subject")
        if (value := msg.get(header))
    ]

    body_text = ""
    body_part = msg.get_body(preferencelist=("plain", "html"))
    if body_part is not None:
        raw_body = body_part.get_content()
        if body_part.get_content_type() == "text/html":
            # get_content() for text/html already returns a decoded str;
            # re-encode so the bytes-based HTML extractor can strip tags the
            # same way the html/epub/odt paths do.
            body_text = _extract_html_text(raw_body.encode("utf-8"))
        else:
            body_text = raw_body

    attachments = (
        [_eml_attachment_label(part) for part in msg.iter_attachments()]
        if msg.is_multipart()
        else []
    )

    sections = []
    if header_lines:
        sections.append("\n".join(header_lines))
    if body_text.strip():
        sections.append(body_text.strip())
    if attachments:
        sections.append(
            f"[{len(attachments)} attachment(s) not extracted: {', '.join(attachments)}]"
        )

    return "\n\n".join(sections)


# ---------------------------------------------------------------------------
# #125 -- EPUB (reuses the HTML extraction path)
# ---------------------------------------------------------------------------

# XML namespaces used by the EPUB Open Container Format / Open Packaging
# Format specs -- needed to find the rootfile (container.xml) and to query
# the manifest/spine (content.opf) with ElementTree's namespace-qualified
# `find`/`findall`.
_EPUB_CONTAINER_NS = "urn:oasis:names:tc:opendocument:xmlns:container"
_EPUB_OPF_NS = {"opf": "http://www.idpf.org/2007/opf"}


def _extract_epub_text(content: bytes, filename: str) -> str:
    """Extract EPUB chapter text in spine (reading) order (#125).

    An EPUB is a ZIP of XHTML documents plus a manifest/spine describing
    their reading order -- this REUSES the existing `_extract_html_text`
    for each chapter rather than writing a second HTML parser, per the
    issue's explicit contract. Spine order (not zip member order) is what
    determines chapter order: `content.opf`'s `<spine>` element is the
    single source of truth for reading sequence, resolved via
    `META-INF/container.xml` -> `content.opf` -> `<manifest>` (id -> href)
    -> `<spine>` (ordered idrefs).

    Failure paths (never a crash, always a non-retryable ApplicationError):
    - Corrupt zip.
    - Missing/unparseable META-INF/container.xml or content.opf.
    - No spine items at all (nothing to determine chapter order from).
    - DRM/encrypted EPUB, signalled by the standard
      META-INF/encryption.xml manifest -- its mere presence means the
      referenced resources cannot be read as plain XHTML, so this is
      checked and rejected BEFORE attempting to parse any chapter.

    All ten failure sites below raise a non-retryable ``ApplicationError``
    rather than a bare exception (#206, same reasoning already applied to
    PDF/JSON/XLSX/PPTX/DOCX in this module -- see `_extract_pdf_text`):
    every one is a pure function of `content`'s bytes, so it fails
    identically on all of Temporal's default 3 retry attempts. Every except
    clause here is scoped to a NARROW exception type
    (`zipfile.BadZipFile`/`KeyError`/`ET.ParseError`), never a broad
    `except Exception` -- so a `MemoryError` from a pathological zip central
    directory or an oversized XML part is never caught here and propagates
    completely unconverted, same discipline `_extract_pdf_text` established
    for its own construction-only wrap (#195 review follow-up). The chapter
    loop below is likewise unwrapped, for the same reason PDF's page loop
    and XLSX's row loop are: a per-chapter failure discovered lazily (e.g.
    `_extract_html_text` OOMing on one pathological chapter) must stay
    retryable, not be swept into `non_retryable=True` by a catch-all around
    the whole loop.
    """
    try:
        zf = zipfile.ZipFile(io.BytesIO(content))
    except zipfile.BadZipFile as e:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' is not a valid EPUB (corrupt zip archive): {e}",
            type="EpubCorruptZip",
            non_retryable=True,
        ) from e

    try:
        names = zf.namelist()
    except zipfile.BadZipFile as e:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' is not a valid EPUB "
            f"(corrupt zip central directory): {e}",
            type="EpubCorruptZip",
            non_retryable=True,
        ) from e

    if "META-INF/encryption.xml" in names:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' is a DRM-protected/encrypted "
            f"EPUB and cannot be extracted.",
            type="EpubDrmProtected",
            non_retryable=True,
        )

    try:
        container_xml = zf.read("META-INF/container.xml")
    except KeyError as e:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' is not a valid EPUB: "
            f"missing META-INF/container.xml",
            type="EpubMissingContainerXml",
            non_retryable=True,
        ) from e

    try:
        container_root = ET.fromstring(container_xml)  # nosec B314 -- see import note
    except ET.ParseError as e:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' has an unparseable META-INF/container.xml: {e}",
            type="EpubUnparseableContainerXml",
            non_retryable=True,
        ) from e

    rootfile = container_root.find(f".//{{{_EPUB_CONTAINER_NS}}}rootfile")
    if rootfile is None or "full-path" not in rootfile.attrib:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' is not a valid EPUB: no rootfile declared",
            type="EpubNoRootfile",
            non_retryable=True,
        )
    opf_path = rootfile.attrib["full-path"]

    try:
        opf_xml = zf.read(opf_path)
    except KeyError as e:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' is not a valid EPUB: "
            f"declared content.opf '{opf_path}' is missing",
            type="EpubMissingContentOpf",
            non_retryable=True,
        ) from e

    try:
        opf_root = ET.fromstring(opf_xml)  # nosec B314 -- see import note
    except ET.ParseError as e:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' has an unparseable content.opf: {e}",
            type="EpubUnparseableContentOpf",
            non_retryable=True,
        ) from e

    # Manifest: item id -> its attributes (href, media-type, properties).
    # `properties` carries EPUB3's "nav"/"cover-image" markers used below to
    # skip navigation and cover items.
    manifest = {
        item.attrib["id"]: item.attrib
        for item in opf_root.iterfind(".//opf:manifest/opf:item", _EPUB_OPF_NS)
        if "id" in item.attrib
    }

    spine_items = opf_root.findall(".//opf:spine/opf:itemref", _EPUB_OPF_NS)
    if not spine_items:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' has no spine -- cannot determine reading order",
            type="EpubNoSpine",
            non_retryable=True,
        )

    # Resolve manifest hrefs relative to content.opf's own directory.
    opf_dir = opf_path.rsplit("/", 1)[0] + "/" if "/" in opf_path else ""

    # Chapters are keyed by their SPINE POSITION (1-based), not by how many
    # have been successfully extracted so far (#125 review blocker 1): a
    # skipped chapter must leave a gap, never shift every later chapter's
    # number down. Renumbering by success-count silently mislabels every
    # chapter after the first miss -- content loss reported as success, with
    # every downstream citation for the rest of the book pointing at the
    # wrong chapter.
    chapters: list[tuple[int, str]] = []
    for position, itemref in enumerate(spine_items, start=1):
        idref = itemref.attrib.get("idref", "")
        item = manifest.get(idref)
        if item is None:
            logger.warning(
                "EPUB spine references a manifest id that does not exist -- "
                "skipping this chapter, not renumbering the rest",
                filename=filename,
                idref=idref,
                spine_position=position,
            )
            continue
        properties = item.get("properties", "").split()
        if "nav" in properties or "cover-image" in properties:
            continue  # nav/cover items skipped per #125 -- not a loss, no warning needed
        if item.get("media-type") not in ("application/xhtml+xml", "text/html"):
            continue  # spine can reference non-markup resources; only chapters are extracted

        # Manifest hrefs are URL references (EPUB OPF spec): spaces and
        # non-ASCII characters are percent-encoded (e.g. "chapter%201.xhtml"
        # for a zip member literally named "chapter 1.xhtml"). Without
        # unquoting, zf.read() misses the real member and this chapter is
        # silently dropped (#125 review blocker 1).
        href = unquote(item.get("href", ""))
        try:
            chapter_bytes = zf.read(opf_dir + href)
        except KeyError:
            logger.warning(
                "EPUB manifest references a zip member that does not exist -- "
                "skipping this chapter, not renumbering the rest",
                filename=filename,
                href=href,
                spine_position=position,
            )
            continue

        chapter_text = _extract_html_text(chapter_bytes).strip()
        if not chapter_text:
            logger.warning(
                "EPUB chapter produced no extractable text -- skipping",
                filename=filename,
                href=href,
                spine_position=position,
            )
            continue

        heading = _epub_chapter_heading(chapter_bytes) or f"Chapter {position}"
        chapters.append((position, f"## {heading}\n\n{chapter_text}"))

    if not chapters:
        raise ApplicationError(
            f"EPUB extraction failed: '{filename}' has a spine but no chapter "
            f"produced any extractable text",
            type="EpubNoExtractableContent",
            non_retryable=True,
        )

    # `chapters` is already in spine (reading) order because the loop above
    # walks `spine_items` in order and only ever appends -- no re-sort
    # needed. The tuple's position element exists for the log lines above,
    # not for reordering here.
    return "\n\n".join(text for _position, text in chapters)


def _epub_chapter_heading(chapter_bytes: bytes) -> str | None:
    """The chapter's own `<title>` or first `<h1>` text, if present --
    preferred over a generic "Chapter N" label so the markdown heading
    carries real information (#125 review blocker 1). Returns None if
    neither exists (or bs4 is unavailable), letting the caller fall back to
    the chapter's spine position.
    """
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return None

    soup = BeautifulSoup(chapter_bytes, "html.parser")
    for tag_name in ("title", "h1"):
        tag = soup.find(tag_name)
        if tag is not None:
            text = tag.get_text(strip=True)
            if text:
                return text
    return None


# ---------------------------------------------------------------------------
# #126 -- RTF and ODT (two distinct formats, two distinct code paths)
# ---------------------------------------------------------------------------


def _extract_rtf_text(content: bytes, filename: str) -> str:
    """Extract text from RTF content via `striprtf` (#126).

    RTF is a control-word format, not XML/ZIP -- kept in its own function
    (not sharing code with `_extract_odt_text` below) so RTF's control-word
    parsing and ODT's zip/XML handling never bleed together.
    `striprtf.rtf_to_text` already handles arbitrarily nested control-word
    groups correctly (that's the whole point of the library), so no bespoke
    nesting logic is needed here.

    Raises:
        ApplicationError (non_retryable=True): striprtf is not installed
            (deterministic per worker/image -- same
            ``MissingExtractionDependency`` reasoning as every other missing
            -library case in this module), `rtf_to_text` cannot parse
            `content` at all, or parsing yields no non-whitespace text. All
            three are deterministic given fixed `content` bytes (#206, same
            fix shape as PDF/JSON/XLSX/PPTX/DOCX above) -- a bare
            `RuntimeError` here was retried 3x by Temporal's default
            RetryPolicy for an outcome already known at attempt 1.

            The try/except around `rtf_to_text` is scoped to ONLY that call
            -- matching `_extract_pdf_text`'s construction-only wrap -- and
            explicitly re-raises `MemoryError` before falling through to the
            broad `except Exception` (#195 review follow-up): a
            load-dependent OOM parsing a pathological/deeply-nested RTF must
            stay retryable, never be reclassified as non-retryable by a
            catch-all.
    """
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as e:
        raise ApplicationError(
            "striprtf not available for RTF extraction",
            type="MissingExtractionDependency",
            non_retryable=True,
        ) from e

    # RTF is a control-word format where non-ASCII text is escaped inline
    # (\'hh hex escapes, \uNNNN unicode control words) rather than encoded
    # directly in the byte stream. latin-1 is used here purely as a
    # lossless byte<->str round trip (maps all 256 byte values 1:1, so it
    # never raises) -- striprtf resolves the ACTUAL characters from those
    # escapes during parsing, not from this decode step.
    raw_text = content.decode("latin-1")

    try:
        # striprtf ships no type stubs, so its return is untyped `Any` to
        # mypy -- str(...) is a no-op at runtime (rtf_to_text always returns
        # a str) but gives the function's own `-> str` signature a real
        # guarantee instead of just trusting an untyped third party.
        text = str(rtf_to_text(raw_text))
    except MemoryError:
        # Load-dependent, not a property of the bytes -- must stay
        # retryable (see the docstring above). Never reclassified as
        # non_retryable.
        raise
    except Exception as e:
        raise ApplicationError(
            f"RTF extraction failed: could not parse '{filename}' "
            f"({type(e).__name__}: {e}). The file may be corrupt or not "
            f"actually RTF despite its declared type.",
            type="RtfParseFailed",
            non_retryable=True,
        ) from e

    if not text.strip():
        raise ApplicationError(
            f"RTF extraction failed: '{filename}' produced no extractable text from RTF content",
            type="RtfNoExtractableText",
            non_retryable=True,
        )
    return text


# ODF (OpenDocument Format) namespaces used by content.xml.
_ODF_OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
_ODF_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"

# Subtrees whose text must NOT be indexed as current document content
# (#126 review blocker 2). content.xml is NOT "just another XML dialect
# with human text inside element bodies" -- it is a real, structured format
# with constructs that carry text an agent must never see as-is:
#   - text:tracked-changes holds DELETED text kept for revision history,
#     not the document's current content.
#   - office:annotation is a PRIVATE reviewer comment (and its dc:creator
#     child is the reviewer's NAME) attached inline to the body, not part
#     of the document itself.
# Indexing either indistinguishably from real content lets a retrieval
# agent cite retracted text as current fact, or surface a private review
# comment (and who wrote it) to anyone with read access to the chunk.
_ODF_EXCLUDED_TAGS = frozenset(
    {
        f"{{{_ODF_TEXT_NS}}}tracked-changes",
        f"{{{_ODF_OFFICE_NS}}}annotation",
    }
)

# Paragraph-level tags: text after one of these ends a line in the
# extracted output, mirroring how a word processor renders them.
_ODF_PARAGRAPH_TAGS = frozenset({f"{{{_ODF_TEXT_NS}}}p", f"{{{_ODF_TEXT_NS}}}h"})

# ODF represents runs of literal spaces/tabs as dedicated elements instead
# of raw whitespace in the XML text (which XML would otherwise collapse) --
# text:s optionally repeats via a text:c count attribute, text:tab is a
# single tab stop. Stripping these tags without substituting real
# whitespace merges adjacent words together; naively treating every tag
# boundary as a newline (the previous bs4-based approach) instead SPLIT one
# paragraph's words across several lines (#126 review blocker 2 secondary
# defect: "Word with spacing" indexed as three lines).
_ODF_SPACE_TAG = f"{{{_ODF_TEXT_NS}}}s"
_ODF_TAB_TAG = f"{{{_ODF_TEXT_NS}}}tab"


def _odf_extract_text(element: ET.Element) -> list[str]:
    """Recursively collect real body text fragments from an ODF element
    tree, in document order -- excluding `_ODF_EXCLUDED_TAGS` subtrees
    entirely and translating whitespace elements to literal whitespace
    (#126 review blocker 2). A paragraph/heading element's fragments end
    with a newline; everything else concatenates inline.
    """
    if element.tag in _ODF_EXCLUDED_TAGS:
        return []  # skip the whole subtree -- tracked changes / annotations

    if element.tag == _ODF_SPACE_TAG:
        # text:c (the repeat count) is itself namespaced -- prefixed XML
        # attributes carry their prefix's namespace, unlike unprefixed ones.
        count = int(element.get(f"{{{_ODF_TEXT_NS}}}c", "1") or "1")
        return [" " * count]
    if element.tag == _ODF_TAB_TAG:
        return ["\t"]

    fragments: list[str] = []
    if element.text:
        fragments.append(element.text)
    for child in element:
        fragments.extend(_odf_extract_text(child))
        if child.tail:
            fragments.append(child.tail)
    if element.tag in _ODF_PARAGRAPH_TAGS:
        fragments.append("\n")
    return fragments


def _extract_odt_text(content: bytes, filename: str) -> str:
    """Extract text from an ODT (OpenDocument Text) `content.xml` (#126).

    ODT is a ZIP container like DOCX/EPUB, but text extraction only needs
    its `content.xml` member -- read via stdlib `zipfile`, then walked with
    stdlib `ElementTree` (no odfpy dependency needed, per the #126
    contract), ODF-STRUCTURE-AWARE rather than a generic tag-strip: real
    ODF documents carry revision-history and reviewer-comment text that must
    never be indexed as current content (see `_odf_extract_text`). Kept in
    its own function (not sharing code with `_extract_rtf_text` above) so
    ODT's zip/XML handling and RTF's control-word parsing never bleed
    together.

    Raises:
        ApplicationError (non_retryable=True): the zip cannot be opened or
            its central directory cannot be read, `content.xml` is missing
            (also the exact signal for a DOCX payload mislabeled as ODT --
            DOCX has `word/document.xml`, not `content.xml`, at its root),
            `content.xml` is unparseable XML, no `office:body/office:text`
            element exists, or the walk yields no text at all. All six are
            deterministic given fixed `content` bytes (#206, same fix shape
            as PDF/JSON/XLSX/PPTX/DOCX/EPUB above) -- a bare `RuntimeError`
            here was retried 3x by Temporal's default RetryPolicy for an
            outcome already known at attempt 1.

            Every except clause below is scoped to a NARROW exception type
            (`zipfile.BadZipFile`/`KeyError`/`ET.ParseError`), never a broad
            `except Exception` -- so a `MemoryError` (e.g. from a
            pathological zip central directory) is never caught here and
            propagates completely unconverted, same discipline
            `_extract_epub_text` above and `_extract_pdf_text` establish.
            `_odf_extract_text`'s recursive body walk below is likewise
            unwrapped, for the same reason.
    """
    try:
        zf = zipfile.ZipFile(io.BytesIO(content))
    except zipfile.BadZipFile as e:
        raise ApplicationError(
            f"ODT extraction failed: '{filename}' is not a valid ODT (corrupt zip archive): {e}",
            type="OdtCorruptZip",
            non_retryable=True,
        ) from e

    try:
        content_xml = zf.read("content.xml")
    except KeyError as e:
        # Also the exact signal for "extension says .odt but the zip is
        # actually a different OOXML/ZIP payload (e.g. DOCX)": a real ODT
        # always has content.xml at its root; DOCX has word/document.xml
        # instead, so this KeyError is a real, actionable contradiction.
        raise ApplicationError(
            f"ODT extraction failed: '{filename}' is not a valid ODT: no "
            f"content.xml found in the archive (the file may be a different "
            f"ZIP-based format, e.g. DOCX, saved with an .odt extension)",
            type="OdtMissingContentXml",
            non_retryable=True,
        ) from e
    except zipfile.BadZipFile as e:
        raise ApplicationError(
            f"ODT extraction failed: '{filename}' is not a valid ODT "
            f"(corrupt zip central directory): {e}",
            type="OdtCorruptZip",
            non_retryable=True,
        ) from e

    try:
        root = ET.fromstring(content_xml)  # nosec B314 -- see import note
    except ET.ParseError as e:
        raise ApplicationError(
            f"ODT extraction failed: '{filename}' has an unparseable content.xml: {e}",
            type="OdtUnparseableContentXml",
            non_retryable=True,
        ) from e

    body = root.find(f".//{{{_ODF_OFFICE_NS}}}body/{{{_ODF_OFFICE_NS}}}text")
    if body is None:
        raise ApplicationError(
            f"ODT extraction failed: '{filename}' is not a valid ODT: no "
            f"office:body/office:text found in content.xml",
            type="OdtMissingBody",
            non_retryable=True,
        )

    raw = "".join(_odf_extract_text(body))
    # Collapse to one non-empty, stripped line per paragraph/heading (mirrors
    # the previous get_text(separator="\n", strip=True) shape callers rely
    # on) -- text:s/text:tab whitespace inserted mid-line above survives
    # this since it isn't a newline.
    text = "\n".join(line.strip() for line in raw.splitlines() if line.strip())

    if not text:
        raise ApplicationError(
            f"ODT extraction failed: '{filename}' produced no extractable text from content.xml",
            type="OdtNoExtractableContent",
            non_retryable=True,
        )
    return text


def _extract_xml_text(content: bytes) -> str:
    """Extract text from XML content (#121) by reusing HTML's tag-stripping
    path verbatim, on the same stdlib `html.parser` BeautifulSoup backend.

    Deliberate parser choice, not an oversight: `html.parser` never resolves
    DTDs or external entities at all, so it carries no XXE / billion-laughs
    entity-expansion attack surface on untrusted uploaded content -- unlike
    a naive `xml.etree.ElementTree.parse` or an `lxml`-backed parser with
    entity resolution enabled. The trade-off is that it is not a strict,
    validating XML parser: malformed XML degrades to best-effort text
    extraction instead of raising, matching YAML/TOML's own "no parse step,
    still searchable" policy (`text_passthrough`) rather than hard-failing
    on a syntax error a human config author made. Attribute VALUES are
    dropped along with their tags -- only element text content survives
    `get_text()`, identical to HTML's existing behavior.

    bs4 warns (``XMLParsedAsHTMLWarning``) when it detects an ``<?xml ...?>``
    declaration being parsed by the HTML backend -- expected and harmless
    given the deliberate choice above, so warnings are suppressed for this
    call rather than left to print on every XML upload. Filtered by
    `simplefilter` (not importing the specific warning class) so this stays
    symmetric with `_extract_html_text`'s own defensive ``ImportError``
    handling -- bs4 IS a core dependency of this service today, but this
    function should degrade the same way that one does, not gain a harder
    unconditional import of its own.
    """
    import warnings

    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        return _extract_html_text(content)


# ---------------------------------------------------------------------------
# #127: SRT / WebVTT subtitle transcript extraction
# ---------------------------------------------------------------------------

# Matches an SRT ("HH:MM:SS,mmm --> HH:MM:SS,mmm") or WebVTT
# ("HH:MM:SS.mmm --> HH:MM:SS.mmm") cue timestamp line. `search`ed against
# each line rather than `match`ed at position 0 so WebVTT's optional cue
# SETTINGS after the second timestamp (e.g. "align:start position:10%")
# don't need special-casing -- they simply trail the match untouched.
_TIMESTAMP_LINE_RE = re.compile(
    r"(\d{2}:\d{2}:\d{2}[.,]\d{3})\s*-->\s*(\d{2}:\d{2}:\d{2}[.,]\d{3})"
)

# A citation marker is re-inserted every Nth cue instead of every cue (which
# would pollute embeddings with a numeric timestamp token per line -- exactly
# the noise stripping cue numbers/timestamps was meant to remove) or never
# (which would make a transcript un-citable -- the whole reason a transcript
# beats a plain summary for an agent). 10 is a citable-but-not-noisy middle
# ground: roughly one marker per 20-60s of realistic dialogue-paced cues,
# close enough for an agent to locate a moment without scrubbing the entire
# file.
_TIMESTAMP_MARKER_EVERY_N_CUES = 10


def _timestamp_to_marker(timestamp: str) -> str:
    """Convert an 'HH:MM:SS[.,]mmm' cue timestamp to a coarse '[t=MM:SS]'
    citation marker (#127's exact proposed format).

    Hours fold into minutes (a 90-minute cue becomes "90:00", not
    "01:30:00") rather than adding a third conditional format -- #127 asks
    specifically for MM:SS, and a >99-minute source is an edge case this
    folding handles correctly without extra branching.
    """
    hours, minutes, seconds = timestamp.replace(",", ".").split(":")
    whole_seconds = seconds.split(".", 1)[0]
    total_minutes = int(hours) * 60 + int(minutes)
    return f"[t={total_minutes:02d}:{int(whole_seconds):02d}]"


def _parse_subtitle_cues(text: str) -> list[tuple[str, str]]:
    """Parse decoded SRT or WebVTT text into ``(start_timestamp, cue_text)``
    pairs, one function for both formats.

    Both share the same cue shape closely enough that one parser handles
    both: a block of lines separated by a blank line, one of which is a
    timestamp line (``_TIMESTAMP_LINE_RE`` accepts SRT's comma or WebVTT's
    dot decimal separator), optionally preceded by a cue number (SRT) or
    identifier (WebVTT), followed by one or more lines of cue text. A block
    with NO timestamp line -- WebVTT's ``WEBVTT`` header, a ``NOTE`` comment
    block, or any other non-cue content -- is simply not a cue and is
    skipped, no special-casing needed for either format's extra bookkeeping.
    """
    # Mixed line endings (bare \r, \r\n, \n in the same file -- real-world
    # subtitle files exported by different tools routinely mix these) must
    # not fragment a single cue's text across a phantom blank-line split, so
    # normalize once up front before splitting into blocks.
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")

    cues: list[tuple[str, str]] = []
    for block in re.split(r"\n\s*\n", normalized):
        lines = [line for line in block.splitlines() if line.strip()]

        timestamp_line_index = None
        match = None
        for index, line in enumerate(lines):
            candidate = _TIMESTAMP_LINE_RE.search(line)
            if candidate is not None:
                timestamp_line_index, match = index, candidate
                break
        # The two are assigned together above, so `match is None` alone is the
        # real condition; the index test is what lets mypy narrow it to int.
        if match is None or timestamp_line_index is None:
            continue  # header / NOTE / cue-number-only / non-cue block

        cue_text = " ".join(lines[timestamp_line_index + 1 :]).strip()
        if cue_text:
            cues.append((match.group(1), cue_text))

    return cues


def _extract_subtitle_text(content: bytes, filename: str) -> str:
    """Extract prose from an SRT or WebVTT transcript (#127).

    Cue numbers and per-cue timestamp lines are stripped entirely -- pure
    retrieval noise, e.g. a ``00:00:04,500 --> 00:00:08,000`` line embeds as
    meaningless numeric tokens -- but a coarse ``[t=MM:SS]`` marker is
    reinserted every `_TIMESTAMP_MARKER_EVERY_N_CUES` cues (see that
    constant's comment) so an agent reading the extracted text can still
    cite roughly WHEN something was said. Timestamps are what makes a
    transcript citable evidence rather than an anonymous paraphrase; losing
    them entirely would make the extraction less useful to the end user
    (an AI agent) than the raw file, and keeping every single one would
    bury the actual spoken content in numeric noise -- this is the
    deliberate middle ground.

    Raises:
        ApplicationError (non_retryable=True): no cue has a recognizable
            timestamp line at all (#127 failure path) -- distinct from the
            generic "empty extraction" guard in `_extract_text_inner`: this
            is specifically "not a valid SRT/WebVTT file", not "a valid one
            that happened to be blank", so the message points at the actual
            problem. Deterministic given fixed `content` bytes (#206, same
            fix shape as PDF/JSON/XLSX/PPTX/DOCX/EPUB/RTF/ODT above) -- a
            bare `RuntimeError` here was retried 3x by Temporal's default
            RetryPolicy for an outcome already known at attempt 1.

            Neither `_decode_text` nor `_parse_subtitle_cues` is wrapped in
            any try/except in this function, so a `MemoryError` from either
            (e.g. charset detection scanning a huge mis-encoded file)
            already propagates completely unconverted.
    """
    text = _decode_text(content)
    cues = _parse_subtitle_cues(text)

    if not cues:
        raise ApplicationError(
            f"Subtitle extraction failed: no subtitle cues with valid timestamps "
            f"found in {filename}. Expected at least one "
            "'HH:MM:SS,mmm --> HH:MM:SS,mmm' (SRT) or 'HH:MM:SS.mmm --> HH:MM:SS.mmm' "
            "(WebVTT) cue.",
            type="SubtitleNoCuesFound",
            non_retryable=True,
        )

    parts: list[str] = []
    for index, (start_timestamp, cue_text) in enumerate(cues):
        if index % _TIMESTAMP_MARKER_EVERY_N_CUES == 0:
            parts.append(_timestamp_to_marker(start_timestamp))
        parts.append(cue_text)
    return " ".join(parts)
